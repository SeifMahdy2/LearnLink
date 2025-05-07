from flask import Flask, request, jsonify, redirect, session, send_file, Response, stream_with_context, make_response
from flask_cors import CORS
import firebase_admin
from firebase_admin import credentials, firestore, auth, storage
import os
from werkzeug.utils import secure_filename
from google_drive_service import GoogleDriveService
from dotenv import load_dotenv
import tempfile
from datetime import datetime, timedelta, date
import uuid
import json
from openai import OpenAI
import markdown
from docxtpl import DocxTemplate
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
import traceback
import time
import re
import hashlib
try:
    from flask_session import Session
except ImportError:
    print("Flask-Session not installed. Use: pip install Flask-Session")
    Session = None

try:
    from serpapi import GoogleSearch
except ImportError:
    print("SerpAPI package not installed. Please run: pip install google-search-results")
    GoogleSearch = None

# Import for the learning style model
from transformers import AutoTokenizer, AutoModelForSequenceClassification
import torch

# Import for document text extraction
import PyPDF2
import docx
from pptx import Presentation

# Load environment variables
dotenv_path = os.path.join(os.path.dirname(__file__), '.env')
if os.path.exists(dotenv_path):
    print(f"Loading .env file from {dotenv_path}")
    load_dotenv(dotenv_path)
else:
    print(f"Warning: .env file not found at {dotenv_path}")
    load_dotenv()  # Try default locations

# Initialize OpenAI API key
openai_api_key = os.getenv("OPENAI_API_KEY")
if not openai_api_key:
    print("WARNING: OPENAI_API_KEY environment variable not set")
    print("Please set the OPENAI_API_KEY environment variable to use OpenAI features")
else:
    print("OpenAI API key loaded successfully")

# Create OpenAI client
client = OpenAI(api_key=openai_api_key)

# Initialize SerpAPI key
serpapi_key = os.getenv("SERPAPI_API_KEY")
if not serpapi_key:
    print("WARNING: SERPAPI_API_KEY environment variable not set")
    print("Please set the SERPAPI_API_KEY environment variable to use image search features")
else:
    print("SerpAPI key loaded successfully")

# Check for credentials file as fallback
def check_credentials_file():
    credentials_file = os.path.join(os.path.dirname(__file__), 'credentials.json')
    if os.path.exists(credentials_file):
        print(f"Found credentials.json file at {credentials_file}")
        try:
            # Read and verify the credentials file
            with open(credentials_file, 'r') as f:
                creds = json.load(f)
                web_creds = creds.get('web', {})
                
                # Set environment variables from credentials file
                if not os.getenv('GOOGLE_DRIVE_CLIENT_ID') and 'client_id' in web_creds:
                    os.environ['GOOGLE_DRIVE_CLIENT_ID'] = web_creds['client_id']
                    print("Set GOOGLE_DRIVE_CLIENT_ID from credentials.json")
                    
                if not os.getenv('GOOGLE_DRIVE_CLIENT_SECRET') and 'client_secret' in web_creds:
                    os.environ['GOOGLE_DRIVE_CLIENT_SECRET'] = web_creds['client_secret']
                    print("Set GOOGLE_DRIVE_CLIENT_SECRET from credentials.json")
                    
                if not os.getenv('GOOGLE_DRIVE_REDIRECT_URI') and 'redirect_uris' in web_creds and web_creds['redirect_uris']:
                    os.environ['GOOGLE_DRIVE_REDIRECT_URI'] = web_creds['redirect_uris'][0]
                    print("Set GOOGLE_DRIVE_REDIRECT_URI from credentials.json")
                    
                return True
        except Exception as e:
            print(f"Error loading credentials.json: {e}")
    return False

# Try to load from credentials.json if env vars are missing
if not (os.getenv('GOOGLE_DRIVE_CLIENT_ID') and 
        os.getenv('GOOGLE_DRIVE_CLIENT_SECRET') and 
        os.getenv('GOOGLE_DRIVE_REDIRECT_URI')):
    check_credentials_file()

# Print environment variables for debugging
if os.getenv('GOOGLE_DRIVE_CLIENT_ID'):
    print("GOOGLE_DRIVE_CLIENT_ID loaded successfully")
else:
    print("WARNING: GOOGLE_DRIVE_CLIENT_ID not found in .env file")
    
if os.getenv('GOOGLE_DRIVE_CLIENT_SECRET'):
    print("GOOGLE_DRIVE_CLIENT_SECRET loaded successfully")
else:
    print("WARNING: GOOGLE_DRIVE_CLIENT_SECRET not found in .env file")
    
if os.getenv('GOOGLE_DRIVE_REDIRECT_URI'):
    print("GOOGLE_DRIVE_REDIRECT_URI loaded successfully")
else:
    print("WARNING: GOOGLE_DRIVE_REDIRECT_URI not found in .env file")

app = Flask(__name__)
# Use a fixed secret key instead of generating a new one each time the server starts
# This will maintain session cookies across server restarts
app.secret_key = os.getenv('FLASK_SECRET_KEY', 'gradproject-fixed-secret-key-for-session-management')
print("Session secret key configured")

# Configure session behavior for more reliable sessions
app.config['SESSION_COOKIE_SECURE'] = False  # Set to True in production with HTTPS
app.config['SESSION_COOKIE_HTTPONLY'] = True
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(days=7)  # Longer session lifetime
app.config['SESSION_TYPE'] = 'filesystem'  # More persistent than cookie-based
app.config['SESSION_USE_SIGNER'] = True
app.config['SESSION_FILE_DIR'] = os.path.join(os.getcwd(), 'flask_session')
# Create session directory if it doesn't exist
if not os.path.exists(app.config['SESSION_FILE_DIR']):
    os.makedirs(app.config['SESSION_FILE_DIR'])
print("Session configuration complete")

# Initialize Flask-Session
if Session is not None:
    Session(app)
    print("Flask-Session initialized")
else:
    print("Warning: Flask-Session not available. Falling back to Flask's default session.")

# Configure CORS for all routes
CORS(app, 
     resources={r"/api/*": {"origins": ["http://localhost:3000"]}},
     supports_credentials=True,
     allow_headers=["Content-Type", "Authorization", "Accept", "X-Requested-With"],
     methods=["GET", "POST", "PUT", "DELETE", "OPTIONS", "PATCH"],
     expose_headers=["Content-Type", "Authorization"])

@app.after_request
def after_request(response):
    # Remove any existing CORS headers to prevent duplication
    response.headers.pop('Access-Control-Allow-Origin', None)
    response.headers.pop('Access-Control-Allow-Headers', None)
    response.headers.pop('Access-Control-Allow-Methods', None)
    response.headers.pop('Access-Control-Allow-Credentials', None)
    
    # Add CORS headers
    response.headers.add('Access-Control-Allow-Origin', 'http://localhost:3000')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization,Accept,X-Requested-With')
    response.headers.add('Access-Control-Allow-Methods', 'GET,POST,PUT,DELETE,OPTIONS,PATCH')
    response.headers.add('Access-Control-Allow-Credentials', 'true')
    return response

# Create temp_uploads directory if it doesn't exist
if not os.path.exists('temp_uploads'):
    os.makedirs('temp_uploads')

# Create files directory for local storage if Google Drive is not available
if not os.path.exists('files'):
    os.makedirs('files')
    # Create subdirectories for different learning styles
    for style in ['visual', 'auditory', 'kinesthetic', 'text']:
        style_dir = os.path.join('files', style)
        if not os.path.exists(style_dir):
            os.makedirs(style_dir)

# Initialize Firebase
firebase_app = None
db = None
bucket = None

try:
    # Check if Firebase credentials are provided in environment variables
    firebase_project_id = os.getenv('FIREBASE_PROJECT_ID')
    firebase_private_key = os.getenv('FIREBASE_PRIVATE_KEY')
    firebase_client_email = os.getenv('FIREBASE_CLIENT_EMAIL')
    
    if firebase_project_id and firebase_private_key and firebase_client_email:
        print("Using Firebase credentials from environment variables")
        # Replace escaped newlines in the private key
        if "\\n" in firebase_private_key:
            firebase_private_key = firebase_private_key.replace("\\n", "\n")
            
        # Create credentials dictionary
        firebase_creds = {
            "type": "service_account",
            "project_id": firebase_project_id,
            "private_key": firebase_private_key,
            "client_email": firebase_client_email
        }
        
        # Initialize Firebase with credentials from environment variables
        cred = credentials.Certificate(firebase_creds)
        firebase_app = firebase_admin.initialize_app(cred, {
            'storageBucket': f"{firebase_project_id}.appspot.com"
        })
        print(f"Firebase app initialized from environment variables with project ID: {firebase_app.project_id}")
        
        # Initialize Firestore
        db = firestore.client()
        print("Firestore client initialized")
        
        # Initialize Storage bucket
        bucket = storage.bucket()
        print(f"Storage bucket initialized: {bucket.name}")
    else:
        # Check if the firebase-adminsdk.json file exists
        firebase_sdk_path = 'firebase-adminsdk.json'
        
        if os.path.exists(firebase_sdk_path):
            print(f"Found Firebase credentials at {firebase_sdk_path}")
            try:
                # Read and verify the credentials file
                with open(firebase_sdk_path, 'r') as f:
                    creds_content = json.load(f)
                    required_fields = ['type', 'project_id', 'private_key_id', 'private_key', 'client_email']
                    missing_fields = [field for field in required_fields if field not in creds_content]
                    
                    if missing_fields:
                        print(f"Error: Firebase credentials file is missing required fields: {', '.join(missing_fields)}")
                        print("Please check the credentials file and ensure it contains all required fields.")
                        print("Continuing without Firebase functionality. Some features may not work.")
                    else:
                        print(f"Using Firebase project: {creds_content.get('project_id')}")
                        
                        # Initialize Firebase with credentials
                        cred = credentials.Certificate(firebase_sdk_path)
                        # Initialize app with storage bucket
                        firebase_app = firebase_admin.initialize_app(cred, {
                            'storageBucket': "grad-project32.firebasestorage.app"
                        })
                        print(f"Firebase app initialized with project ID: {firebase_app.project_id}")
                        
                        # Initialize Firestore
                        db = firestore.client()
                        print("Firestore client initialized")
                        
                        # Initialize Storage bucket
                        bucket = storage.bucket(name="grad-project32.firebasestorage.app")
                        print(f"Storage bucket initialized: {bucket.name}")
                        
                        # Test the bucket
                        try:
                            bucket_metadata = bucket.get_iam_policy()
                            print("Storage bucket access verified")
                        except Exception as bucket_error:
                            print(f"Warning: Could not access bucket metadata: {str(bucket_error)}")
                            print("This may indicate permission issues with the bucket.")
                        
                        print("Firebase initialized successfully")
            except json.JSONDecodeError:
                print(f"Error: Firebase credentials file is not valid JSON: {firebase_sdk_path}")
                print("Firebase functionality will be disabled. Some features may not work.")
            except Exception as init_error:
                print(f"Error initializing Firebase: {str(init_error)}")
                traceback.print_exc()
                print("Firebase functionality will be disabled. Some features may not work.")
        else:
            print(f"Warning: Firebase credentials file not found at {firebase_sdk_path}")
            print("Firebase functionality will be disabled. Some features may not work.")
except Exception as e:
    print(f"Critical error during Firebase setup: {str(e)}")
    traceback.print_exc()
    print("Firebase functionality will be disabled. Some features may not work.")

# Store drive service instances for each user
drive_services = {}

# Initialize the learning style model (lazy loading to save resources)
learning_style_model = None
learning_style_tokenizer = None
learning_style_labels = ['Auditory', 'Kinesthetic', 'Reading/Writing', 'Visual']

def load_learning_style_model():
    """Load the model and tokenizer if not already loaded"""
    global learning_style_model, learning_style_tokenizer
    
    if learning_style_model is None:
        try:
            model_name = "pushpikaLiyanagama/student-learning-style-identify"
            learning_style_tokenizer = AutoTokenizer.from_pretrained(model_name)
            learning_style_model = AutoModelForSequenceClassification.from_pretrained(model_name)
            print("Learning style model loaded successfully")
        except Exception as e:
            print(f"Error loading learning style model: {e}")
            return False
    return True

def predict_learning_style(text):
    """Predict learning style from text using the Hugging Face model"""
    if not load_learning_style_model():
        return None
    
    try:
        inputs = learning_style_tokenizer(text, return_tensors="pt", padding=True, truncation=True)
        outputs = learning_style_model(**inputs)
        logits = outputs.logits
        predicted_class = torch.argmax(logits, dim=1).item()
        return learning_style_labels[predicted_class]
    except Exception as e:
        print(f"Error predicting learning style: {e}")
        return None

def get_drive_service():
    """Get the drive service for the current user."""
    user_email = session.get('user_email')
    
    # If no user is in session, use a default service
    if not user_email:
        # Use a default service for anonymous uploads
        if 'default' not in drive_services:
            drive_services['default'] = GoogleDriveService()
        return drive_services['default']
        
    # Otherwise use the user-specific service
    if user_email not in drive_services:
        drive_services[user_email] = GoogleDriveService(user_email)
    return drive_services[user_email]

@app.route("/api/check-user", methods=["POST"])
def check_user():
    """Check if a user exists in Firestore"""
    try:
        data = request.json
        if not data:
            print("Error in check_user: No JSON data provided")
            return jsonify({"error": "No JSON data provided"}), 400
            
        email = data.get("email")
        
        if not email:
            print("Error in check_user: Email is required")
            return jsonify({"error": "Email is required"}), 400
            
        print(f"Checking if user exists with email: {email}")
            
        # Check if user exists
        users_ref = db.collection("users")
        query = users_ref.where("email", "==", email).limit(1)
        
        # Execute query and check if any documents match
        results = list(query.stream())
        
        exists = len(results) > 0
        user_id = results[0].id if results else None
        
        print(f"User check result - exists: {exists}, user_id: {user_id}")
        
        return jsonify({
            "exists": exists,
            "user_id": user_id
        })
    except Exception as e:
        print(f"Error in check_user: {str(e)}")
        return jsonify({"error": f"Failed to check user: {str(e)}"}), 500

@app.route("/signup", methods=["POST"])
def signup():
    try:
        data = request.json
        if not data:
            return jsonify({"error": "Invalid request: No JSON data provided"}), 400
            
        email = data.get("email")
        name = data.get("name")
        auth_provider = data.get("authProvider", "email")
        uid = data.get("uid")
        photo_url = data.get("photoURL", "")
        created_at = data.get("createdAt", datetime.now().isoformat())

        if not email or not name:
            return jsonify({"error": "Email and name are required"}), 400

        # Check if user already exists
        users_ref = db.collection("users")
        existing_users = users_ref.where("email", "==", email).stream()

        for user in existing_users:
            return jsonify({"error": "User already exists in Firestore"}), 400

        # Store user in Firestore
        user_ref = db.collection("users").document()
        user_data = {
            "name": name,
            "email": email,
            "authProvider": auth_provider,
            "createdAt": created_at
        }
        
        # Add optional fields if they exist
        if uid:
            user_data["uid"] = uid
        if photo_url:
            user_data["photoURL"] = photo_url

        user_ref.set(user_data)

        return jsonify({
            "message": "User added to Firestore", 
            "user_id": user_ref.id
        }), 201
    except Exception as e:
        print(f"Error in signup route: {str(e)}")
        return jsonify({"error": f"Failed to create user: {str(e)}"}), 500

def allowed_file(filename):
    """Check if the file extension is allowed"""
    ALLOWED_EXTENSIONS = {'pdf', 'doc', 'docx', 'txt', 'ppt', 'pptx'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Upload a file to Firebase Storage and save metadata in Firestore"""
    try:
        print("\n=== Starting file upload process ===")
        
        # 1. Check request
        print("1. Checking request...")
        if 'file' not in request.files:
            print("Error: No file part in request")
            return jsonify({"error": "No file part"}), 400
        
        file = request.files['file']
        if file.filename == '':
            print("Error: No selected file")
            return jsonify({"error": "No selected file"}), 400
            
        # Check if file type is allowed
        if not allowed_file(file.filename):
            print(f"Error: File type not allowed for {file.filename}")
            return jsonify({
                "error": "File type not allowed. Please upload only PDF, DOC, DOCX, TXT, or PowerPoint files."
            }), 400
        
        user_id = request.form.get('userId')
        if not user_id:
            print("Error: No user ID provided")
            return jsonify({"error": "User ID is required"}), 400
        
        learning_style = request.form.get('learningStyle', 'visual')
        print(f"Request validated - File: {file.filename}, User: {user_id}, Style: {learning_style}")
        
        # 2. Generate file ID and prepare paths
        file_id = str(uuid.uuid4())
        filename = secure_filename(file.filename)
        print(f"Generated file ID: {file_id}")
        print(f"Secured filename: {filename}")
        
        # 3. Save to temp file
        print("\n2. Saving to temporary file...")
        temp_dir = os.path.join(os.getcwd(), 'temp_uploads')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            print(f"Created temp_uploads directory at {temp_dir}")
        
        temp_path = os.path.join(temp_dir, filename)
        file.save(temp_path)
        print(f"File saved temporarily at: {temp_path}")
        
        # 4. Debug Firebase Storage state
        print("\n3. Checking Firebase Storage initialization...")
        
        # Get global variables
        global bucket, firebase_app, db
        
        # Verify Firebase is initialized
        if not firebase_app:
            print("Error: Firebase app not initialized")
            return jsonify({"error": "Firebase app not initialized"}), 500
            
        if not db:
            print("Error: Firestore client not initialized")
            return jsonify({"error": "Firestore client not initialized"}), 500
            
        if not bucket:
            print("Error: Firebase Storage bucket not initialized")
            # Try to re-initialize the bucket
            try:
                print("Attempting to re-initialize storage bucket...")
                bucket = storage.bucket(name="grad-project32.firebasestorage.app")
                print(f"Re-initialized storage bucket: {bucket.name}")
            except Exception as bucket_error:
                print(f"Failed to re-initialize bucket: {str(bucket_error)}")
                return jsonify({"error": "Firebase Storage not available. Could not initialize bucket."}), 500

        # Print bucket name for verification
        try:
            bucket_name = bucket.name
            print(f"Using Firebase Storage bucket: {bucket_name}")
        except Exception as e:
            print(f"Error accessing bucket name: {str(e)}")
            return jsonify({"error": f"Firebase Storage bucket error: {str(e)}"}), 500

        # 5. Create folder structure and upload to Firebase Storage
        print("\n4. Creating folder structure in Firebase Storage...")
        # Define storage path with folder structure
        storage_path = f"users/{user_id}/{file_id}/{filename}"
        print(f"Storage path: {storage_path}")

        try:
            # Create an empty placeholder file to ensure folder structure exists
            folder_path = f"users/{user_id}/{file_id}/"
            placeholder = bucket.blob(folder_path + ".placeholder")
            placeholder.upload_from_string("")
            print(f"Created folder structure: {folder_path}")
            
            # Now create the actual file blob
            blob = bucket.blob(storage_path)
            print("Created blob reference")
            
            print(f"Uploading file from {temp_path}...")
            with open(temp_path, 'rb') as file_obj:
                blob.upload_from_file(file_obj)
            print("File uploaded to Firebase Storage successfully")
            
            # 6. Make file public and get URL
            print("\n5. Making file public...")
            blob.make_public()
            download_url = blob.public_url
            print(f"File made public. URL: {download_url}")
            
            # 7. Save metadata to Firestore
            print("\n6. Saving metadata to Firestore...")
            file_data = {
                "name": filename,
                "type": file.content_type,
                "size": os.path.getsize(temp_path),
                "storagePath": storage_path,
                "downloadUrl": download_url,
                "userId": user_id,
                "learningStyle": learning_style,
                "status": "uploaded",
                "processed": False,
                "createdAt": firestore.SERVER_TIMESTAMP
            }
            
            # Save to Firestore
            file_ref = db.collection("files").document(file_id)
            file_ref.set(file_data)
            print(f"Metadata saved to Firestore. Document ID: {file_id}")
            
            # 8. Clean up
            print("\n7. Cleaning up...")
            if os.path.exists(temp_path):
                os.remove(temp_path)
                print(f"Temporary file removed: {temp_path}")
            
            print("\n=== Upload process completed successfully ===")
            return jsonify({
                "success": True,
                "fileId": file_id,
                "downloadUrl": download_url,
                "message": "File uploaded successfully"
            })
        except Exception as storage_error:
            print("🔥 STORAGE ERROR 🔥")
            print(f"Firebase Storage error: {str(storage_error)}")
            traceback.print_exc()
            
            # More specific error message
            error_message = str(storage_error)
            if "The specified bucket does not exist" in error_message:
                error_message = "The Firebase Storage bucket does not exist. Please check your Firebase configuration."
            elif "Permission denied" in error_message:
                error_message = "Permission denied when accessing Firebase Storage. Please check your Firebase access rules."
            return jsonify({
                "success": False,
                "error": f"Firebase Storage error: {error_message}"
            }), 500
            
    except Exception as e:
        print(f"\nError during upload: {str(e)}")
        print("Detailed error trace:")
        traceback.print_exc()
        
        # Clean up temp file if it exists
        if 'temp_path' in locals() and os.path.exists(temp_path):
            os.remove(temp_path)
            print(f"Cleaned up temporary file after error: {temp_path}")
        
        return jsonify({
            "success": False,
            "error": f"Failed to upload file: {str(e)}"
        }), 500

@app.route('/api/drive/upload', methods=['POST'])
def upload_to_drive():
    if 'file' not in request.files:
        return jsonify({'success': False, 'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'success': False, 'error': 'No file selected'}), 400
        
    # Check if file type is allowed
    if not allowed_file(file.filename):
        return jsonify({
            'success': False,
            'error': 'File type not allowed. Please upload only PDF, DOC, DOCX, TXT, or PowerPoint files.'
        }), 400

    # Get learning style from form data if available
    learning_style = request.form.get('learningStyle', 'visual')
    
    # Check if we should use local storage (for development or when Google Drive auth fails)
    use_local_storage = request.form.get('useLocalStorage', 'false').lower() == 'true'

    try:
        # Create temp directory if it doesn't exist
        if not os.path.exists('temp_uploads'):
            os.makedirs('temp_uploads')
            
        filename = secure_filename(file.filename)
        temp_path = os.path.join('temp_uploads', filename)
        
        # Save file temporarily
        file.save(temp_path)
        
        # If explicitly requested local storage, don't try Google Drive
        if use_local_storage:
            # Save to local storage directory
            local_file_dir = os.path.join('files', learning_style)
            if not os.path.exists(local_file_dir):
                os.makedirs(local_file_dir)
                
            local_path = os.path.join(local_file_dir, filename)
            import shutil
            shutil.copy(temp_path, local_path)
            
            # Clean up temporary file
            if os.path.exists(temp_path):
                os.remove(temp_path)
                
            return jsonify({
                'success': True,
                'warning': 'Using local storage as requested.',
                'file_path': local_path,
                'name': filename,
                'learning_style': learning_style
            })
        
        # Try Google Drive upload
        drive_service = get_drive_service()
        
        # Check if service is authenticated
        if not drive_service.is_authenticated():
            # If using the frontend directly, return an authentication error
            if user_agent := request.headers.get('User-Agent'):
                if 'mozilla' in user_agent.lower():
                    auth_url = drive_service.get_authorization_url()
                    return jsonify({
                        'success': False, 
                        'error': 'Authentication required',
                        'auth_url': auth_url
                    }), 401
            
            # For direct API calls, attempt to store in local filesystem as fallback
            # This is for development/testing only
            local_file_dir = os.path.join('files', learning_style)
            if not os.path.exists(local_file_dir):
                os.makedirs(local_file_dir)
                
            local_path = os.path.join(local_file_dir, filename)
            import shutil
            shutil.copy(temp_path, local_path)
            
            # Return success with a warning about local storage
            if os.path.exists(temp_path):
                os.remove(temp_path)
                
            return jsonify({
                'success': True,
                'warning': 'Not authenticated with Google Drive. File stored locally.',
                'file_path': local_path,
                'learning_style': learning_style
            })
        
        # Include learning style in metadata
        result = drive_service.upload_file(temp_path, filename, learning_style)
        
        # Clean up temporary file
        if os.path.exists(temp_path):
            os.remove(temp_path)
        
        return jsonify(result)
    except Exception as e:
        # Clean up temp file if it exists
        temp_path = os.path.join('temp_uploads', secure_filename(file.filename))
        if os.path.exists(temp_path):
            os.remove(temp_path)
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/drive/download/<file_id>', methods=['GET'])
def download_from_drive(file_id):
    try:
        # Get custom filename from query parameter
        custom_filename = request.args.get('filename')
        
        # Get file metadata first
        drive_service = get_drive_service()
        file_metadata = drive_service.get_file_metadata(file_id)
        if not file_metadata['success']:
            return jsonify(file_metadata), 400

        # Create temp directory if it doesn't exist
        if not os.path.exists('temp_uploads'):
            os.makedirs('temp_uploads')

        # Generate temp file path using original name for correct extension
        temp_path = os.path.join('temp_uploads', file_metadata['name'])
        
        # Download the file
        result = drive_service.download_file(file_id, temp_path)
        
        if result['success']:
            try:
                # If custom filename is provided, use it with original extension
                if custom_filename:
                    original_ext = os.path.splitext(file_metadata['name'])[1]
                    download_name = f"{custom_filename}{original_ext}"
                else:
                    download_name = file_metadata['name']

                # Send the file to the client
                response = send_file(
                    temp_path,
                    as_attachment=True,
                    download_name=download_name
                )
                # Clean up the temp file after sending
                @response.call_on_close
                def cleanup():
                    if os.path.exists(temp_path):
                        os.remove(temp_path)
                return response
            except Exception as e:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
                raise e
        else:
            return jsonify(result), 400
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/drive/files', methods=['GET'])
def list_drive_files():
    try:
        page_size = request.args.get('pageSize', default=10, type=int)
        drive_service = get_drive_service()
        result = drive_service.list_files(page_size=page_size)
        return jsonify(result)
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/drive/auth/status')
def auth_status():
    """Check if user is authenticated."""
    user_email = session.get('user_email')
    if user_email and user_email in drive_services:
        drive_service = drive_services[user_email]
        return jsonify({
            'isAuthenticated': drive_service.is_authenticated(),
            'user': {'email': user_email}
        })
    return jsonify({'isAuthenticated': False, 'user': None})

@app.route('/api/drive/auth/url')
def auth_url():
    """Get the authorization URL."""
    drive_service = GoogleDriveService()
    url = drive_service.get_authorization_url()
    return jsonify({'url': url})

@app.route('/oauth2callback')
def oauth2callback():
    """Handle the OAuth callback."""
    try:
        code = request.args.get('code')
        if not code:
            return jsonify({'success': False, 'error': 'No code provided'})

        drive_service = GoogleDriveService()
        user_email = drive_service.handle_oauth2_callback(code)
        
        if user_email:
            session['user_email'] = user_email
            drive_services[user_email] = drive_service
            return """
                <script>
                    window.opener.postMessage('authentication_successful', '*');
                    window.close();
                </script>
            """
        return jsonify({'success': False, 'error': 'Authentication failed'})
    except Exception as e:
        print(f"OAuth callback error: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/drive/auth/logout', methods=['POST'])
def logout():
    """Handle user logout."""
    user_email = session.get('user_email')
    if user_email in drive_services:
        del drive_services[user_email]
    session.clear()
    return jsonify({'success': True})

@app.route('/api/drive/delete/<file_id>', methods=['DELETE'])
def delete_file(file_id):
    """Delete a file from Google Drive"""
    try:
        drive_service = get_drive_service()
        result = drive_service.delete_file(file_id)
        return jsonify(result)
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/predict-learning-style', methods=['POST'])
def predict_style():
    """Endpoint to predict learning style from text"""
    data = request.json
    text = data.get("text", "")
    
    if not text:
        return jsonify({"success": False, "error": "No text provided"}), 400
    
    try:
        style = predict_learning_style(text)
        
        if style:
            # If user is authenticated, store their learning style in Firebase
            user_email = session.get('user_email')
            if user_email:
                users_ref = db.collection("users")
                user_query = users_ref.where("email", "==", user_email).limit(1)
                user_docs = list(user_query.stream())
                
                if user_docs:
                    user_id = user_docs[0].id
                    # Update user document with learning style
                    db.collection("users").document(user_id).update({
                        "learning_style": style,
                        "learning_style_updated_at": firestore.SERVER_TIMESTAMP
                    })
                    
            return jsonify({"success": True, "learning_style": style})
        else:
            return jsonify({"success": False, "error": "Could not predict learning style"}), 500
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/store-learning-style', methods=['POST'])
def store_learning_style():
    """Endpoint to store user's learning style in Firestore"""
    data = request.json
    learning_style = data.get("learningStyle")
    learning_style_details = data.get("learningStyleDetails", {})
    email = data.get("email")
    uid = data.get("uid")
    
    if not learning_style or not email:
        return jsonify({"success": False, "error": "Learning style and email are required"}), 400
    
    try:
        # Find user in Firestore
        users_ref = db.collection("users")
        user_query = users_ref.where("email", "==", email).limit(1)
        user_docs = list(user_query.stream())
        
        if user_docs:
            # Update existing user
            user_id = user_docs[0].id
            db.collection("users").document(user_id).update({
                "learningStyle": learning_style,
                "learningStyleDetails": learning_style_details,
                "learningStyleUpdatedAt": firestore.SERVER_TIMESTAMP
            })
            return jsonify({"success": True, "message": "Learning style updated"})
        else:
            # User not found, let's create one
            user_data = {
                "email": email,
                "learningStyle": learning_style,
                "learningStyleDetails": learning_style_details,
                "learningStyleUpdatedAt": firestore.SERVER_TIMESTAMP,
                "createdAt": firestore.SERVER_TIMESTAMP,
            }
            
            # Add optional fields
            if uid:
                user_data["uid"] = uid
            if data.get("name"):
                user_data["name"] = data.get("name")
                
            new_user_ref = db.collection("users").document()
            new_user_ref.set(user_data)
            return jsonify({"success": True, "message": "New user created with learning style"})
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/store-quiz-score', methods=['POST'])
def store_quiz_score():
    """Endpoint to store user's quiz score in Firestore"""
    data = request.json
    email = data.get("email")
    uid = data.get("uid")
    quiz_id = data.get("quizId")
    document_id = data.get("documentId")
    quiz_name = data.get("quizName", "Quiz")
    score = data.get("score")
    correct_answers = data.get("correctAnswers")
    total_questions = data.get("totalQuestions")
    percentage = data.get("percentage")
    
    # Validate required fields
    if not email or not score or not quiz_id:
        return jsonify({"success": False, "error": "Email, score, and quiz ID are required"}), 400
    
    try:
        # Find user in Firestore
        users_ref = db.collection("users")
        user_query = users_ref.where("email", "==", email).limit(1)
        user_docs = list(user_query.stream())
        
        user_id = None
        
        if user_docs:
            # User exists
            user_id = user_docs[0].id
            user_ref = db.collection("users").document(user_id)
            
            # Check if quizzes field exists and create it if it doesn't
            user_data = user_docs[0].to_dict()
            
            if "quizScores" not in user_data:
                user_ref.update({"quizScores": []})
            
            # Add new quiz score
            quiz_score_data = {
                "quizId": quiz_id,
                "quizName": quiz_name,
                "score": score,
                "correctAnswers": correct_answers,
                "totalQuestions": total_questions,
                "percentage": percentage,
                "timestamp": firestore.SERVER_TIMESTAMP
            }
            
            # Add document ID if available
            if document_id:
                quiz_score_data["documentId"] = document_id
            
            # Update user document with new quiz score
            user_ref.update({
                "quizScores": firestore.ArrayUnion([quiz_score_data])
            })
            
            return jsonify({"success": True, "message": "Quiz score saved successfully"})
        else:
            # User not found, let's create one with quiz scores
            user_data = {
                "email": email,
                "quizScores": [{
                    "quizId": quiz_id,
                    "quizName": quiz_name,
                    "score": score,
                    "correctAnswers": correct_answers,
                    "totalQuestions": total_questions,
                    "percentage": percentage,
                    "timestamp": firestore.SERVER_TIMESTAMP
                }],
                "createdAt": firestore.SERVER_TIMESTAMP,
            }
            
            # Add optional fields
            if uid:
                user_data["uid"] = uid
            if data.get("name"):
                user_data["name"] = data.get("name")
                
            # Add document ID if available
            if document_id:
                user_data["quizScores"][0]["documentId"] = document_id
                
            new_user_ref = db.collection("users").document()
            new_user_ref.set(user_data)
            
            return jsonify({"success": True, "message": "New user created with quiz score"})
        
    except Exception as e:
        print(f"Error saving quiz score: {str(e)}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/user/quiz-scores', methods=['POST'])
def get_user_quiz_scores():
    """Endpoint to get user's quiz scores from Firestore"""
    data = request.json
    email = data.get("email")
    uid = data.get("uid")
    
    # Validate required fields
    if not email:
        return jsonify({"success": False, "error": "Email is required"}), 400
    
    try:
        # Find user in Firestore
        users_ref = db.collection("users")
        user_query = users_ref.where("email", "==", email).limit(1)
        user_docs = list(user_query.stream())
        
        if user_docs:
            # User exists
            user_data = user_docs[0].to_dict()
            
            # Get quiz scores
            quiz_scores = user_data.get("quizScores", [])
            
            # Convert timestamps to strings for JSON serialization
            for score in quiz_scores:
                if "timestamp" in score and score["timestamp"]:
                    if hasattr(score["timestamp"], "seconds"):
                        score["timestamp"] = score["timestamp"].strftime("%Y-%m-%d %H:%M:%S")
            
            return jsonify({
                "success": True, 
                "quizScores": quiz_scores
            })
        else:
            # User not found
            return jsonify({
                "success": False, 
                "error": "User not found",
                "quizScores": []
            }), 404
        
    except Exception as e:
        print(f"Error getting quiz scores: {str(e)}")
        return jsonify({
            "success": False, 
            "error": str(e),
            "quizScores": []
        }), 500

@app.route('/api/files/local', methods=['GET'])
def list_local_files():
    """List files stored locally for development/testing when Google Drive is not available."""
    try:
        files_list = []
        learning_style = request.args.get('learningStyle', None)
        
        base_dir = 'files'
        if learning_style:
            # Look in specific style directory if specified
            style_dir = os.path.join(base_dir, learning_style)
            if os.path.exists(style_dir):
                for filename in os.listdir(style_dir):
                    file_path = os.path.join(style_dir, filename)
                    if os.path.isfile(file_path):
                        file_size = os.path.getsize(file_path)
                        modified_time = os.path.getmtime(file_path)
                        files_list.append({
                            'name': filename,
                            'path': file_path,
                            'size': file_size,
                            'modified': modified_time,
                            'learning_style': learning_style
                        })
        else:
            # Look in all style directories
            for style in os.listdir(base_dir):
                style_dir = os.path.join(base_dir, style)
                if os.path.isdir(style_dir):
                    for filename in os.listdir(style_dir):
                        file_path = os.path.join(style_dir, filename)
                        if os.path.isfile(file_path):
                            file_size = os.path.getsize(file_path)
                            modified_time = os.path.getmtime(file_path)
                            files_list.append({
                                'name': filename,
                                'path': file_path,
                                'size': file_size,
                                'modified': modified_time,
                                'learning_style': style
                            })
        
        return jsonify({
            'success': True,
            'files': files_list,
            'note': 'These files are stored locally and not in Google Drive'
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/files/user/<user_id>', methods=['GET'])
def get_user_files(user_id):
    """Get all files for a specific user from Firestore"""
    try:
        # Query Firestore for files that belong to this user
        files_ref = db.collection('files').where('userId', '==', user_id).order_by('createdAt', direction=firestore.Query.DESCENDING)
        files_docs = files_ref.get()
        
        # Convert to list of dictionaries
        files = []
        for doc in files_docs:
            file_data = doc.to_dict()
            file_data['id'] = doc.id
            files.append(file_data)
            
        return jsonify({
            'success': True,
            'files': files
        })
    except Exception as e:
        print(f"Error getting files for user {user_id}: {str(e)}")
        return jsonify({
            'success': False,
            'error': f"Failed to retrieve files: {str(e)}"
        }), 500

@app.route('/api/files/<file_id>', methods=['DELETE'])
def delete_file_metadata(file_id):
    """Delete file metadata from Firestore"""
    try:
        print(f"Attempting to delete file metadata with ID: {file_id}")
        
        # Check for required parameters
        if not file_id:
            return jsonify({"error": "File ID is required"}), 400
        
        # Get the file document from Firestore
        file_ref = db.collection("files").document(file_id)
        file_doc = file_ref.get()
        
        if not file_doc.exists:
            return jsonify({"error": "File not found"}), 404
        
        # Delete from Firestore
        file_ref.delete()
        print(f"Deleted file metadata from Firestore with ID: {file_id}")
        
        return jsonify({
            "success": True,
            "message": "File metadata deleted successfully"
        })
        
    except Exception as e:
        print(f"Error deleting file: {str(e)}")
        return jsonify({
            "success": False,
            "error": f"Failed to delete file: {str(e)}"
        }), 500

@app.route('/api/files', methods=['GET'])
def get_files():
    """Get files for a specific user from Firestore"""
    try:
        user_id = request.args.get('userId')
        if not user_id:
            return jsonify({"error": "User ID is required"}), 400
        print(f"Fetching files for user: {user_id}")
        
        files = []
        
        # Query files in Firestore
        files_ref = db.collection("files")
        query = files_ref.where("userId", "==", user_id).order_by("createdAt", direction=firestore.Query.DESCENDING)
        
        for doc in query.stream():
            file_data = doc.to_dict()
            file_data['id'] = doc.id  # Add document ID
            files.append(file_data)
        
        print(f"Found {len(files)} files for user {user_id} in Firestore")
        
        return jsonify({
            "success": True,
            "files": files,
            "source": "firebase"
        })
        
    except Exception as e:
        print(f"Error getting files: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Failed to get files: {str(e)}"}), 500

@app.route('/api/files/<file_id>/process', methods=['POST'])
def process_document(file_id):
    """
    Process a document using AI. Fetches the document from Firebase Storage,
    processes its content, and updates its status in Firestore.
    """
    try:
        print(f"Processing document with ID: {file_id}")
        
        # Get the document metadata from Firestore
        doc_ref = db.collection("files").document(file_id)
        doc = doc_ref.get()
        
        if not doc.exists:
            print(f"Document with ID {file_id} not found")
            return jsonify({"error": "Document not found"}), 404
        
        doc_data = doc.to_dict()
        
        # Update status to processing
        doc_ref.update({
            "processingStatus": "processing",
            "updatedAt": datetime.now()
        })
        
        # Get storage path from metadata
        storage_path = doc_data.get("storagePath")
        if not storage_path:
            print(f"Storage path not found in document metadata")
            doc_ref.update({
                "processingStatus": "failed",
                "processingError": "Storage path not found in document metadata",
                "updatedAt": datetime.now()
            })
            return jsonify({"error": "Storage path not found"}), 400
        
        # Create temp dir if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
        
        # Download file from Firebase Storage
        temp_path = os.path.join(temp_dir, f"{file_id}_{os.path.basename(storage_path)}")
        blob = bucket.blob(storage_path)
        
        try:
            blob.download_to_filename(temp_path)
            print(f"Downloaded file for processing to: {temp_path}")
        except Exception as e:
            error_msg = f"Error downloading file: {str(e)}"
            print(error_msg)
            doc_ref.update({
                "processingStatus": "failed",
                "processingError": error_msg,
                "updatedAt": datetime.now()
            })
            return jsonify({"error": error_msg}), 500
        
        # Extract text from the document
        try:
            # This is where you would implement text extraction based on file type
            # For now, we'll simulate this with a placeholder
            extracted_text = "This is placeholder text extracted from the document."
            print(f"Text extracted from document: {len(extracted_text)} characters")
            
            # Generate content with AI based on learning style
            learning_style = doc_data.get("learningStyle", "visual")
            
            # This is where you would implement AI content generation
            # For now, we'll simulate with a placeholder
            processed_content = f"Processed content for {learning_style} learning style: {extracted_text}"
            
            # Create a processed version of the document
            processed_filename = f"processed_{file_id}.txt"
            processed_path = os.path.join(temp_dir, processed_filename)
            
            with open(processed_path, "w") as f:
                f.write(processed_content)
            
            # Upload processed document to Firebase Storage
            processed_storage_path = f"processed/{doc_data.get('userId', 'unknown')}/{processed_filename}"
            processed_blob = bucket.blob(processed_storage_path)
            processed_blob.upload_from_filename(processed_path)
            
            # Make the processed file publicly accessible
            processed_blob.make_public()
            processed_download_url = processed_blob.public_url
            
            # Update document metadata with processing results
            doc_ref.update({
                "processingStatus": "completed",
                "processedContent": processed_content,
                "processedStoragePath": processed_storage_path,
                "processedDownloadUrl": processed_download_url,
                "updatedAt": datetime.now()
            })
            
            # Clean up temporary files
            if os.path.exists(temp_path):
                os.remove(temp_path)
            if os.path.exists(processed_path):
                os.remove(processed_path)
            
            return jsonify({
                "success": True,
                "message": "Document processed successfully",
                "fileId": file_id,
                "downloadUrl": processed_download_url
            })
            
        except Exception as e:
            error_msg = f"Error processing document: {str(e)}"
            print(error_msg)
            traceback.print_exc()
            doc_ref.update({
                "processingStatus": "failed",
                "processingError": error_msg,
                "updatedAt": datetime.now()
            })
            return jsonify({"error": error_msg}), 500
        
    except Exception as e:
        print(f"Error processing document: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Error processing document: {str(e)}"}), 500

@app.route('/api/download/<file_id>/<format>', methods=['GET'])
def download_processed_document(file_id, format):
    """Download a processed document from Firebase Storage"""
    try:
        print(f"Attempting to download processed document with ID: {file_id} in format: {format}")
        
        # Get the file metadata from Firestore
        file_doc = db.collection("files").document(file_id).get()
        
        if not file_doc.exists:
            print(f"File with ID {file_id} not found in Firestore")
            return jsonify({"error": "File not found"}), 404
        
        file_data = file_doc.to_dict()
        
        # Check if the file has been processed
        if file_data.get("processingStatus") != "completed":
            print(f"File with ID {file_id} has not been processed yet or processing failed")
            return jsonify({
                "error": "File has not been processed yet or processing failed", 
                "status": file_data.get("processingStatus", "unknown")
            }), 400
        
        # Get processed file path from metadata
        processed_storage_path = file_data.get("processedStoragePath")
        if not processed_storage_path:
            print(f"Processed storage path not found in file metadata")
            return jsonify({"error": "Processed file path not found"}), 400
        
        # Create temp dir if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_downloads')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
        
        # Generate a temp file path
        temp_path = os.path.join(temp_dir, f"processed_{file_id}.{format}")
        
        # Download processed file from Firebase Storage
        blob = bucket.blob(processed_storage_path)
        
        try:
            blob.download_to_filename(temp_path)
            print(f"Downloaded processed file to: {temp_path}")
        except Exception as e:
            print(f"Error downloading processed file: {str(e)}")
            return jsonify({"error": f"Error downloading processed file: {str(e)}"}), 500
        
        # Generate appropriate filename
        original_filename = file_data.get("name", f"document_{file_id}")
        base_name = os.path.splitext(original_filename)[0]
        download_name = f"{base_name}_processed.{format}"
        
        # Serve the file with cleanup
        response = send_file(
            temp_path,
            as_attachment=True,
            download_name=download_name
        )
        
        # Clean up after sending
        @response.call_on_close
        def cleanup():
            if os.path.exists(temp_path):
                os.remove(temp_path)
                print(f"Cleaned up temporary file: {temp_path}")
        
        return response
    
    except Exception as e:
        print(f"Error serving processed file: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Error serving processed file: {str(e)}"}), 500

def create_preflight_response():
    """Create a preflight response for CORS requests"""
    response = jsonify({})
    response.headers.add('Access-Control-Allow-Origin', 'http://localhost:3000')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization,Accept,X-Requested-With')
    response.headers.add('Access-Control-Allow-Methods', 'GET,POST,PUT,DELETE,OPTIONS,PATCH')
    response.headers.add('Access-Control-Allow-Credentials', 'true')
    return response

@app.route('/api/files/<file_id>/style', methods=['OPTIONS', 'PATCH'])
def update_learning_style(file_id):
    # Handle preflight request
    if request.method == 'OPTIONS':
        return create_preflight_response()

    try:
        # Get new learning style from request body
        data = request.get_json()
        if not data or 'learningStyle' not in data:
            return jsonify({'success': False, 'error': 'Learning style not provided'}), 400
        
        new_style = data['learningStyle']
        valid_styles = ['visual', 'auditory', 'reading_writing', 'kinesthetic']
        if new_style not in valid_styles:
            return jsonify({'success': False, 'error': f'Invalid learning style. Must be one of: {", ".join(valid_styles)}'}), 400
        
        # Update document in Firestore
        print(f"Updating file {file_id} learning style to: {new_style}")
        
        # Update in Firestore
        db.collection('files').document(file_id).update({
            'learningStyle': new_style,
            'updatedAt': firestore.SERVER_TIMESTAMP
        })
        
        return jsonify({
            'success': True,
            'message': 'Learning style updated successfully',
            'file_id': file_id,
            'learning_style': new_style
        }), 200
        
    except Exception as e:
        print(f"Error in update_learning_style: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/files/<file_id>', methods=['PATCH'])
def update_file(file_id):
    """Update file properties in Firestore"""
    try:
        print(f"Attempting to update file with ID: {file_id}")
        
        # Check for required parameters
        if not file_id:
            return jsonify({"success": False, "error": "File ID is required"}), 400
            
        # Get request data
        data = request.get_json()
        if not data:
            return jsonify({"success": False, "error": "No update data provided"}), 400
            
        # Get the file document from Firestore
        file_ref = db.collection("files").document(file_id)
        file_doc = file_ref.get()
        
        if not file_doc.exists:
            return jsonify({"success": False, "error": "File not found"}), 404
            
        # Check what fields we're updating
        update_data = {}
        
        # Handle learning style update
        if 'learningStyle' in data:
            # Validate learning style
            valid_styles = ['visual', 'auditory', 'reading_writing', 'kinesthetic']
            if data['learningStyle'] not in valid_styles:
                return jsonify({"success": False, "error": f"Invalid learning style. Must be one of: {', '.join(valid_styles)}"}), 400
                
            update_data['learningStyle'] = data['learningStyle']
            print(f"Updating learning style to: {data['learningStyle']}")
            
        # Add other updateable fields as needed
        
        # If nothing to update, return error
        if not update_data:
            return jsonify({"success": False, "error": "No valid update fields provided"}), 400
            
        # Add updated timestamp
        update_data['updatedAt'] = firestore.SERVER_TIMESTAMP
            
        # Update the document
        file_ref.update(update_data)
        print(f"Updated file {file_id} with data: {update_data}")
        
        # Get the updated document
        updated_doc = file_ref.get()
        updated_data = updated_doc.to_dict()
        
        return jsonify({
            "success": True,
            "message": "File updated successfully",
            "file": updated_data
        }), 200
            
    except Exception as e:
        print(f"Error updating file: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

# Set up OpenAI API key
openai_api_key = os.getenv("OPENAI_API_KEY")
if not openai_api_key:
    print("WARNING: OPENAI_API_KEY environment variable not set")
    print("Please set the OPENAI_API_KEY environment variable to use OpenAI features")
else:
    print("OpenAI API key loaded successfully")

# Create OpenAI client
client = OpenAI(api_key=openai_api_key)

# Initialize SerpAPI key
serpapi_key = os.getenv("SERPAPI_API_KEY")
if not serpapi_key:
    print("WARNING: SERPAPI_API_KEY environment variable not set")
    print("Please set the SERPAPI_API_KEY environment variable to use image search features")
else:
    print("SerpAPI key loaded successfully")

# Function to extract text from documents based on file type
def extract_text_from_document(file_path):
    """Extract text from various document formats"""
    file_extension = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_extension == '.pdf':
            return extract_text_from_pdf(file_path)
        elif file_extension in ['.docx', '.doc']:
            return extract_text_from_docx(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            return extract_text_from_pptx(file_path)
        elif file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        else:
            return f"Unsupported file format: {file_extension}"
    except Exception as e:
        print(f"Error extracting text from document: {e}")
        return f"Error extracting text: {str(e)}"

def extract_text_from_pdf(file_path):
    """Extract text from PDF files"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return f"Error extracting PDF text: {str(e)}"

def extract_text_from_docx(file_path):
    """Extract text from DOCX files"""
    try:
        doc = docx.Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return f"Error extracting DOCX text: {str(e)}"

def extract_text_from_pptx(file_path):
    """Extract text from PPTX files"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return f"Error extracting PPTX text: {str(e)}"

def generate_reading_writing_content(document_text):
    """
    Generate content optimized for reading/writing learning style using OpenAI
    """
    if not openai_api_key:
        # Return a mock response for development without API key
        print("Using mock response for reading/writing content (No API key)")
        return {
            "title": "Reading/Writing Learning Materials",
            "description": "This content has been optimized for reading/writing learners.",
            "elements": [
                {
                    "type": "text",
                    "content": "# Main Concepts\n\n" + 
                              "The document covers the following key areas:\n\n" +
                              "1. Introduction to the topic\n" +
                              "2. Core principles and methodologies\n" +
                              "3. Practical applications\n\n" +
                              "## Detailed Notes\n\n" +
                              "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Ut at ex ac libero pellentesque molestie. Donec interdum ante et nisi dignissim, in consectetur tellus porttitor.\n\n" +
                              "### Key Points:\n" +
                              "* Point 1: Important information about the topic\n" +
                              "* Point 2: Critical analysis of the concept\n" +
                              "* Point 3: Practical implementation strategies\n\n" +
                              "## Summary\n\n" +
                              "This section summarizes the main ideas presented in the document and provides a coherent overview of the material.",
                    "caption": "Structured Notes"
                }
            ]
        }
        
    try:
        # Truncate the document_text if it's too long for a single API call
        truncated_text = document_text[:7000] if len(document_text) > 7000 else document_text
        
        # First, get an overview and main topics from the document
        overview_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert educator specializing in identifying main topics and creating structured outlines. Extract the main topics from this document and create a brief outline."},
                {"role": "user", "content": f"Extract the 3-6 main topics from this document and provide a brief overview:\n\n{truncated_text[:1500]}"}
            ],
            max_tokens=500,
            temperature=0.5
        )
        
        overview_text = overview_response.choices[0].message.content
        
        # Divide the document into chunks of approximately 3000 characters
        chunk_size = 3000
        chunks = []
        for i in range(0, len(document_text), chunk_size):
            chunks.append(document_text[i:i + chunk_size])
        
        # Limit to 5 chunks maximum to avoid excessive API calls
        chunks = chunks[:5]
        
        # Process each chunk separately
        processed_chunks = []
        
        for i, chunk in enumerate(chunks):
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": """You are an expert educator specializing in creating detailed educational content. Your goal is to explain concepts thoroughly while maintaining engagement and clarity. Create content that teaches concepts as if explaining to someone learning about them for the first time.

Follow these guidelines:
1. Start each major section with a brief overview paragraph that introduces the concept
2. Break down complex topics into clear explanations with:
   - Detailed paragraphs that thoroughly explain each concept
   - Bullet points for key features, components, or steps
   - Examples and real-world applications to reinforce understanding
3. Use "Understanding X" sections to provide deeper insights into important concepts
4. Include "Key Concept" boxes (using > for blockquotes) to highlight crucial information
5. Add "For Example" sections to demonstrate practical applications

Format using:
- Clear headings (# for main topics, ## for subtopics, ### for specific concepts)
- Paragraphs for detailed explanations
- Bullet points for features and components
- Blockquotes (>) for key concept boxes
- Numbered lists for steps and processes"""},
                    {"role": "user", "content": f"You're processing part {i+1} of a larger document. Create educational content that explains the concepts in this chunk, ensuring it flows well as part of a larger study guide:\n\n{chunk}"}
                ],
                max_tokens=1500,
                temperature=0.7
            )
            
            # Extract the formatted text from the API response
            formatted_chunk = response.choices[0].message.content
            processed_chunks.append(formatted_chunk)
        
        # Combine the overview with the processed chunks
        combined_content = f"# Document Study Guide\n\n## Overview\n\n{overview_text}\n\n"
        for i, chunk in enumerate(processed_chunks):
            combined_content += f"\n## Part {i+1}\n\n{chunk}\n\n"
            
        combined_content += "\n## Study Tips\n\n* Review each section thoroughly before proceeding to the next\n* Create your own notes based on the key points\n* Try to explain these concepts in your own words\n* Practice applying these concepts to real-world scenarios"
            
        # Structure the response
        return {
            "title": "Reading/Writing Learning Materials",
            "description": "This content has been optimized for reading/writing learners with structured notes, clear headings, and organized points.",
            "elements": [
                {
                    "type": "text",
                    "content": combined_content,
                    "caption": "Structured Study Guide"
                }
            ]
        }
        
    except Exception as e:
        print(f"Error generating reading/writing content: {e}")
        # Return a basic response with the error
        return {
            "title": "Reading/Writing Learning Materials (Error)",
            "description": "There was an error processing this document.",
            "elements": [
                {
                    "type": "text",
                    "content": f"Error generating reading/writing content: {str(e)}\n\nOriginal text (partial):\n\n{document_text[:500]}...",
                    "caption": "Error Processing Document"
                }
            ]
        }

def create_docx_study_guide(content, title="Study Guide"):
    """
    Create a well-formatted DOCX document from the AI-generated content
    - This version creates a DOCX with exactly the same content as the web version
    """
    print(f"Creating DOCX study guide with exact same content as web")
    
    doc = Document()
    
    # Set up styles
    styles = doc.styles
    
    # Title style
    title_style = styles.add_style('CustomTitle', WD_STYLE_TYPE.PARAGRAPH)
    title_font = title_style.font
    title_font.size = Pt(24)
    title_font.bold = True
    title_font.color.rgb = RGBColor(48, 84, 150)
    
    # Heading 1 style
    h1_style = styles.add_style('CustomH1', WD_STYLE_TYPE.PARAGRAPH)
    h1_font = h1_style.font
    h1_font.size = Pt(18)
    h1_font.bold = True
    h1_font.color.rgb = RGBColor(94, 110, 220)
    
    # Heading 2 style
    h2_style = styles.add_style('CustomH2', WD_STYLE_TYPE.PARAGRAPH)
    h2_font = h2_style.font
    h2_font.size = Pt(16)
    h2_font.bold = True
    h2_font.color.rgb = RGBColor(63, 81, 181)
    
    # Heading 3 style
    h3_style = styles.add_style('CustomH3', WD_STYLE_TYPE.PARAGRAPH)
    h3_font = h3_style.font
    h3_font.size = Pt(14)
    h3_font.bold = True
    h3_font.color.rgb = RGBColor(94, 110, 220)
    
    # Normal text style
    normal_style = styles.add_style('CustomNormal', WD_STYLE_TYPE.PARAGRAPH)
    normal_font = normal_style.font
    normal_font.size = Pt(11)
    
    # Add title with consistent text across all platforms
    doc.add_paragraph("Reading/Writing Learning Materials", 'CustomTitle').alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()  # Add some space
    
    # Process each line of content
    current_list = []
    in_list = False
    
    for line in content.split('\n'):
        line = line.strip()
        if not line:
            if in_list and current_list:
                # Add the current list to document
                p = doc.add_paragraph(style='CustomNormal')
                for item in current_list:
                    p = doc.add_paragraph(item, style='CustomNormal')
                    p.style = 'List Bullet'
                current_list = []
                in_list = False
            doc.add_paragraph()  # Add empty paragraph for spacing
            continue
        
        # Handle headers
        if line.startswith('# '):
            if in_list: # End any current list
                p = doc.add_paragraph(style='CustomNormal')
                for item in current_list:
                    p = doc.add_paragraph(item, style='CustomNormal')
                    p.style = 'List Bullet'
                current_list = []
                in_list = False
            doc.add_paragraph(line[2:], 'CustomH1')
        elif line.startswith('## '):
            if in_list:
                p = doc.add_paragraph(style='CustomNormal')
                for item in current_list:
                    p = doc.add_paragraph(item, style='CustomNormal')
                    p.style = 'List Bullet'
                current_list = []
                in_list = False
            doc.add_paragraph(line[3:], 'CustomH2')
        elif line.startswith('### '):
            if in_list:
                p = doc.add_paragraph(style='CustomNormal')
                for item in current_list:
                    p = doc.add_paragraph(item, style='CustomNormal')
                    p.style = 'List Bullet'
                current_list = []
                in_list = False
            doc.add_paragraph(line[4:], 'CustomH3')
        # Handle bullet points
        elif line.startswith('* ') or line.startswith('- '):
            in_list = True
            current_list.append(line[2:])
        # Handle numbered lists
        elif line[0].isdigit() and '. ' in line:
            number, text = line.split('. ', 1)
            p = doc.add_paragraph(style='CustomNormal')
            p.style = 'List Number'
            p.text = text
        # Regular paragraph
        else:
            if in_list:
                p = doc.add_paragraph(style='CustomNormal')
                for item in current_list:
                    p = doc.add_paragraph(item, style='CustomNormal')
                    p.style = 'List Bullet'
                current_list = []
                in_list = False
            doc.add_paragraph(line, 'CustomNormal')
    
    # Save to BytesIO
    docx_file = BytesIO()
    doc.save(docx_file)
    docx_file.seek(0)
    return docx_file

def create_pdf_study_guide(content, title="Study Guide"):
    """
    Create a well-formatted PDF document from the AI-generated content
    - This version creates a PDF with exactly the same content as the web version
    """
    print(f"Creating new PDF study guide with exact same content as web")
    
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
    
    # Get sample stylesheet and modify it
    styles = getSampleStyleSheet()
    
    # Create custom styles
    styles.add(ParagraphStyle(
        name='CustomTitle',
        parent=styles['Title'],
        fontSize=24,
        spaceAfter=30,
        textColor=colors.HexColor('#304E96')
    ))
    
    styles.add(ParagraphStyle(
        name='CustomH1',
        parent=styles['Heading1'],
        fontSize=18,
        spaceAfter=16,
        textColor=colors.HexColor('#5E6EDC')
    ))
    
    styles.add(ParagraphStyle(
        name='CustomH2',
        parent=styles['Heading2'],
        fontSize=16,
        spaceAfter=14,
        textColor=colors.HexColor('#3F51B5')
    ))
    
    styles.add(ParagraphStyle(
        name='CustomH3',
        parent=styles['Heading3'],
        fontSize=14,
        spaceAfter=12,
        textColor=colors.HexColor('#5E6EDC')
    ))
    
    styles.add(ParagraphStyle(
        name='CustomBody',
        parent=styles['Normal'],
        fontSize=11,
        spaceAfter=12
    ))
    
    # Build the PDF content
    elements = []
    
    # Add title - using a consistent title for all documents
    elements.append(Paragraph("Reading/Writing Learning Materials", styles['CustomTitle']))
    elements.append(Spacer(1, 12))
    
    current_list = []
    in_list = False
    
    for line in content.split('\n'):
        line = line.strip()
        if not line:
            if in_list and current_list:
                elements.append(ListFlowable(
                    [ListItem(Paragraph(item, styles['CustomBody'])) for item in current_list],
                    bulletType='bullet'
                ))
                current_list = []
                in_list = False
            elements.append(Spacer(1, 12))
            continue
        
        # Handle headers
        if line.startswith('# '):
            if in_list:
                elements.append(ListFlowable(
                    [ListItem(Paragraph(item, styles['CustomBody'])) for item in current_list],
                    bulletType='bullet'
                ))
                current_list = []
                in_list = False
            elements.append(Paragraph(line[2:], styles['CustomH1']))
        elif line.startswith('## '):
            if in_list:
                elements.append(ListFlowable(
                    [ListItem(Paragraph(item, styles['CustomBody'])) for item in current_list],
                    bulletType='bullet'
                ))
                current_list = []
                in_list = False
            elements.append(Paragraph(line[3:], styles['CustomH2']))
        elif line.startswith('### '):
            if in_list:
                elements.append(ListFlowable(
                    [ListItem(Paragraph(item, styles['CustomBody'])) for item in current_list],
                    bulletType='bullet'
                ))
                current_list = []
                in_list = False
            elements.append(Paragraph(line[4:], styles['CustomH3']))
        # Handle bullet points
        elif line.startswith('* ') or line.startswith('- '):
            in_list = True
            current_list.append(line[2:])
        # Handle numbered lists
        elif line[0].isdigit() and '. ' in line:
            number, text = line.split('. ', 1)
            elements.append(Paragraph(f"{number}. {text}", styles['CustomBody']))
        # Regular paragraph
        else:
            if in_list:
                elements.append(ListFlowable(
                    [ListItem(Paragraph(item, styles['CustomBody'])) for item in current_list],
                    bulletType='bullet'
                ))
                current_list = []
                in_list = False
            elements.append(Paragraph(line, styles['CustomBody']))
    
    # Build the PDF
    doc.build(elements)
    buffer.seek(0)
    return buffer

def test_disabled_storage():
    """Inform the user that storage is disabled"""
    return jsonify({
        "message": "File storage is disabled. Files are only registered in the database but not physically stored.",
        "status": "disabled"
    })

def test_storage_connection():
    """Test connection to Firestore"""
    results = {
        "firebase_app": False,
        "firestore": False,
        "storage": "disabled",
        "errors": []
    }
    
    try:
        # Test Firebase App
        app = firebase_admin.get_app()
        results["firebase_app"] = True
        
        # Test Firestore
        try:
            db.collection("test").limit(1).get()
            results["firestore"] = True
        except Exception as e:
            results["errors"].append(f"Firestore error: {str(e)}")
    
    except Exception as e:
        results["errors"].append(f"Firebase app error: {str(e)}")
    
    return results

@app.route('/api/storage-status', methods=['GET'])
def storage_status():
    """Get status of storage and Firebase services"""
    status = test_storage_connection()
    all_services_ok = all([
        status["firebase_app"],
        status["firestore"]
    ])
    
    return jsonify({
        "success": all_services_ok,
        "status": status,
        "storage_message": "File storage is disabled"
    })

@app.route('/api/storage-test', methods=['GET'])
def test_disabled_storage():
    """Inform the user that storage is disabled"""
    return jsonify({
        "message": "File storage is disabled. Files are only registered in the database but not physically stored.",
        "status": "disabled"
    })

@app.route('/api/files/<file_id>/stream', methods=['GET'])
@app.route('/api/files/<file_id>/download', methods=['GET'])
def file_not_available(file_id):
    """Inform user that files are not stored"""
    return jsonify({
        "error": "File storage is disabled. Files are only registered but not stored on the server.",
        "message": "To enable file storage, please configure a storage solution."
    }), 404

@app.route('/api/files/<file_id>/download', methods=['GET'])
def download_file(file_id):
    """Download a file from Firebase Storage"""
    try:
        print(f"Attempting to download file with ID: {file_id}")
        
        # Get the file metadata from Firestore
        file_doc = db.collection("files").document(file_id).get()
        
        if not file_doc.exists:
            print(f"File with ID {file_id} not found in Firestore")
            return jsonify({"error": "File not found"}), 404
        
        file_data = file_doc.to_dict()
        
        # Get the storage path from metadata
        storage_path = file_data.get("storagePath")
        if not storage_path:
            print(f"Storage path not found in file metadata")
            return jsonify({"error": "Storage path not found"}), 400
        
        # Create temp dir if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_downloads')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
        
        # Generate a temp file path
        temp_filename = f"{file_id}_{os.path.basename(storage_path)}"
        temp_path = os.path.join(temp_dir, temp_filename)
        
        # Download file from Firebase Storage
        blob = bucket.blob(storage_path)
        blob.download_to_filename(temp_path)
        print(f"Downloaded file to: {temp_path}")
        
        # Get the original filename
        filename = file_data.get("name", os.path.basename(storage_path))
        
        # Serve the file with cleanup
        response = send_file(
            temp_path,
            as_attachment=True,
            download_name=filename
        )
        
        # Clean up after sending (using response.call_on_close)
        @response.call_on_close
        def cleanup():
            if os.path.exists(temp_path):
                os.remove(temp_path)
                print(f"Cleaned up temporary file: {temp_path}")
        
        return response
    
    except Exception as e:
        print(f"Error serving file: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Error serving file: {str(e)}"}), 500

@app.route('/api/files/<file_id>/url', methods=['GET'])
def get_file_url(file_id):
    """Get the download URL for a file"""
    try:
        # Get the file metadata from Firestore
        file_doc = db.collection("files").document(file_id).get()
        
        if not file_doc.exists:
            return jsonify({"error": "File not found"}), 404
        
        file_data = file_doc.to_dict()
        
        # Return the download URL from metadata
        if "downloadUrl" in file_data:
            return jsonify({
                "fileId": file_id,
                "downloadUrl": file_data["downloadUrl"],
                "name": file_data.get("name", "unknown")
            })
        else:
            # If URL not in metadata, generate a download endpoint URL
            return jsonify({
                "fileId": file_id,
                "downloadUrl": f"/api/files/{file_id}/download",
                "name": file_data.get("name", "unknown")
            })
        
    except Exception as e:
        print(f"Error getting file URL: {str(e)}")
        return jsonify({"error": f"Error getting file URL: {str(e)}"}), 500

@app.route('/api/files/<file_id>/generate-summary', methods=['POST'])
def generate_summary(file_id):
    """Generate a summary of the document using OpenAI."""
    try:
        # Get the document from Firestore
        doc_ref = db.collection("files").document(file_id)
        doc = doc_ref.get()
        
        if not doc.exists:
            return jsonify({"error": "Document not found"}), 404
            
        doc_data = doc.to_dict()
        storage_path = doc_data.get("storagePath")
        
        if not storage_path:
            return jsonify({"error": "Storage path not found"}), 400
            
        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            
        # Download file
        temp_path = os.path.join(temp_dir, f"{file_id}_{os.path.basename(storage_path)}")
        blob = bucket.blob(storage_path)
        blob.download_to_filename(temp_path)
        
        # Extract text from document
        document_text = extract_text_from_document(temp_path)
        
        # Clean up temp file
        if os.path.exists(temp_path):
            os.remove(temp_path)
            
        # Generate summary using OpenAI
        max_chars = 14000
        truncated_text = document_text[:max_chars] if len(document_text) > max_chars else document_text
        
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert at summarizing documents. Create a comprehensive summary that includes: 1) Main points and key ideas 2) Important details and examples 3) Conclusions or findings. Format the summary with clear sections and bullet points where appropriate."},
                {"role": "user", "content": f"Please summarize this document:\n\n{truncated_text}"}
            ],
            max_tokens=1000
        )
        
        summary = response.choices[0].message.content
        
        # Store the summary in Firestore
        doc_ref.update({
            "summary": summary,
            "summaryGeneratedAt": firestore.SERVER_TIMESTAMP
        })
        
        return jsonify({
            "success": True,
            "summary": summary
        })
        
    except Exception as e:
        print(f"Error generating summary: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/files/<file_id>/generate-quiz', methods=['POST'])
def generate_quiz(file_id):
    try:
        data = request.get_json()
        quiz_type = data.get('quiz_type', 'multiple_choice')
        
        # Get document from Firestore
        doc_ref = db.collection('files').document(file_id)
        doc = doc_ref.get()
        
        if not doc.exists:
            return jsonify({'success': False, 'error': 'Document not found'}), 404
            
        doc_data = doc.to_dict()
        storage_path = doc_data.get('storagePath')
        
        if not storage_path:
            return jsonify({'success': False, 'error': 'Storage path not found'}), 400

        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            
        # Download file
        temp_path = os.path.join(temp_dir, f"{file_id}_{os.path.basename(storage_path)}")
        try:
            blob = bucket.blob(storage_path)
            blob.download_to_filename(temp_path)
            print(f"Downloaded file to: {temp_path}")
        except Exception as e:
            print(f"Error downloading file: {str(e)}")
            return jsonify({'success': False, 'error': f'Error downloading file: {str(e)}'}), 500

        # Extract text from document
        try:
            content = extract_text_from_document(temp_path)
            print(f"Extracted text length: {len(content)}")
            
            # Clean up temp file
            if os.path.exists(temp_path):
                os.remove(temp_path)
                print("Cleaned up temporary file")

            if not content or len(content.strip()) == 0:
                return jsonify({'success': False, 'error': 'Could not extract content from document'}), 400

        except Exception as e:
            print(f"Error extracting text: {str(e)}")
            if os.path.exists(temp_path):
                os.remove(temp_path)
            return jsonify({'success': False, 'error': f'Error extracting text: {str(e)}'}), 500

        # Generate quiz based on type
        if quiz_type == 'multiple_choice':
            system_prompt = """Create 10 multiple choice questions based on the content. Format your response as a JSON array of questions. Each question should have:
            - A clear question text
            - Four options (A, B, C, D)
            - The correct answer letter
            - A brief explanation
            
            Example format:
            {
                "questions": [
                    {
                        "question": "What is the main topic discussed?",
                        "options": [
                            "A) First option",
                            "B) Second option",
                            "C) Third option",
                            "D) Fourth option"
                        ],
                        "correct_answer": "A",
                        "explanation": "This is correct because..."
                    }
                ]
            }"""
        else:
            system_prompt = """Create 5 fill-in-the-blank questions based on the content. For each question:
            1. Take a direct quote or sentence from the content
            2. Replace a key term or concept with _____ (5 underscores)
            3. The blank should be a single word or short phrase that is clearly defined in the content
            4. Include the surrounding context to make it clear what should go in the blank
            
            Format your response as a JSON object with this structure:
            {
                "questions": [
                    {
                        "text_before_blank": "In machine learning, a",
                        "text_after_blank": "neural network is designed to process grid-like data such as images.",
                        "correct_answer": "convolutional",
                        "alternative_answers": ["CNN", "convolution", "convolutional neural network"],
                        "required_keywords": ["conv"],
                        "explanation": "Convolutional neural networks (CNNs) are specialized for processing grid-like data, particularly images, by using convolutional layers to detect patterns and features.",
                        "context": "This type of neural network is specifically used for image processing and computer vision tasks."
                    }
                ]
            }
            
            Make sure each question:
            1. Has a clear blank indicated by _____ (5 underscores)
            2. Has enough context to determine the answer
            3. Has a specific, unambiguous answer
            4. Includes relevant synonyms or acceptable variations of the answer"""

        # Truncate content if too long
        max_chars = 14000
        truncated_content = content[:max_chars] if len(content) > max_chars else content

        # Call OpenAI API
        try:
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"Create quiz questions based on this content:\n\n{truncated_content}"}
                ],
                response_format={ "type": "json_object" },
                temperature=0.7,
                max_tokens=2000
            )

            # Parse the response
            quiz_data = json.loads(response.choices[0].message.content)
            
            # Store the quiz in Firestore
            doc_ref.update({
                f"quiz_{quiz_type}": quiz_data,
                f"quiz_{quiz_type}_generatedAt": firestore.SERVER_TIMESTAMP
            })
            
            return jsonify({
                "success": True,
                "questions": quiz_data["questions"]
            })
            
        except Exception as e:
            print(f"Error generating quiz with OpenAI: {str(e)}")
            return jsonify({'success': False, 'error': f'Error generating quiz: {str(e)}'}), 500
        
    except Exception as e:
        print(f"Error in quiz generation: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/files/<file_id>/process-reading-writing', methods=['OPTIONS', 'POST'])
def process_reading_writing(file_id):
    """Process a document for reading/writing learners by generating a structured study guide."""
    # Handle preflight request
    if request.method == 'OPTIONS':
        return create_preflight_response()
        
    try:
        print(f"Processing reading/writing content for file ID: {file_id}")
        
        # Get the document from Firestore
        doc_ref = db.collection("files").document(file_id)
        doc = doc_ref.get()
        
        if not doc.exists:
            print(f"Document with ID {file_id} not found")
            return jsonify({"error": "Document not found"}), 404
            
        doc_data = doc.to_dict()
        storage_path = doc_data.get("storagePath")
        
        if not storage_path:
            print(f"Storage path not found for document {file_id}")
            return jsonify({"error": "Storage path not found"}), 400
            
        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            
        # Download file
        temp_path = os.path.join(temp_dir, f"{file_id}_{os.path.basename(storage_path)}")
        blob = bucket.blob(storage_path)
        blob.download_to_filename(temp_path)
        
        # Extract text from document
        document_text = extract_text_from_document(temp_path)
        
        # Clean up temp file
        if os.path.exists(temp_path):
            os.remove(temp_path)
            
        # Generate reading/writing optimized content
        content = generate_reading_writing_content(document_text)
        
        # Create both DOCX and PDF versions
        docx_file = create_docx_study_guide(content["elements"][0]["content"], content["title"])
        pdf_file = create_pdf_study_guide(content["elements"][0]["content"], content["title"])
        
        # Upload processed documents to Firebase Storage
        docx_filename = f"reading_writing_{file_id}.docx"
        pdf_filename = f"reading_writing_{file_id}.pdf"
        
        docx_storage_path = f"processed/{doc_data.get('userId', 'unknown')}/{docx_filename}"
        pdf_storage_path = f"processed/{doc_data.get('userId', 'unknown')}/{pdf_filename}"
        
        # Upload DOCX
        docx_blob = bucket.blob(docx_storage_path)
        docx_blob.upload_from_file(docx_file, content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        docx_blob.make_public()
        
        # Upload PDF
        pdf_blob = bucket.blob(pdf_storage_path)
        pdf_blob.upload_from_file(pdf_file, content_type='application/pdf')
        pdf_blob.make_public()
        
        # Store the processed content in Firestore
        doc_ref.update({
            "readingWritingContent": content,
            "readingWritingDocxUrl": docx_blob.public_url,
            "readingWritingPdfUrl": pdf_blob.public_url,
            "readingWritingGeneratedAt": firestore.SERVER_TIMESTAMP
        })
        
        return jsonify({
            "success": True,
            "content": content,
            "docxUrl": docx_blob.public_url,
            "pdfUrl": pdf_blob.public_url
        })
        
    except Exception as e:
        print(f"Error processing reading/writing content: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

def generate_auditory_content(document_text):
    """
    Generate content optimized for auditory learners using OpenAI
    """
    if not openai_api_key:
        print("Using mock response for auditory content (No API key)")
        return {
            "title": "Audio Learning Materials",
            "description": "This content has been optimized for auditory learners with spoken explanations and examples.",
            "elements": [
                {
                    "type": "text",
                    "content": "Imagine you're learning about this topic in a conversation with a friendly tutor...",
                    "caption": "Spoken Explanation"
                }
            ]
        }
    
    try:
        # Truncate text if too long (OpenAI has token limits)
        max_chars = 14000
        truncated_text = document_text[:max_chars] if len(document_text) > max_chars else document_text
        
        # Generate spoken-friendly content first
        chat_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": """You are an AI tutor specialized in creating spoken-friendly content for auditory learners. Transform educational content into engaging, conversational explanations that sound natural when spoken aloud."""},
                {"role": "user", "content": f"Transform this content into a spoken-friendly, engaging explanation that would be easy to understand when read aloud:\n\n{truncated_text}"}
            ],
            max_tokens=2000,
            temperature=0.7
        )
        
        # Extract the formatted text
        formatted_text = chat_response.choices[0].message.content
        
        # Generate audio using OpenAI's Text-to-Speech
        speech_response = client.audio.speech.create(
            model="tts-1",
            voice="alloy",
            input=formatted_text
        )
        
        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            
        # Save the audio file temporarily
        temp_audio_path = os.path.join(temp_dir, f"audio_explanation_{int(time.time())}.mp3")
        speech_response.stream_to_file(temp_audio_path)
        
        # Upload to Firebase Storage
        audio_storage_path = f"audio/audio_explanation_{int(time.time())}.mp3"
        audio_blob = bucket.blob(audio_storage_path)
        audio_blob.upload_from_filename(temp_audio_path, content_type='audio/mpeg')
        audio_blob.make_public()
        
        # Clean up temp file
        if os.path.exists(temp_audio_path):
            os.remove(temp_audio_path)
        
        # Structure the response
        return {
            "title": "Audio Learning Materials",
            "description": "This content has been optimized for auditory learners with spoken explanations and examples.",
            "elements": [
                {
                    "type": "text",
                    "content": formatted_text,
                    "caption": "Spoken Explanation"
                }
            ],
            "audioUrl": audio_blob.public_url
        }
        
    except Exception as e:
        print(f"Error generating auditory content: {e}")
        return {
            "title": "Audio Learning Materials (Error)",
            "description": "There was an error processing this document.",
            "elements": [
                {
                    "type": "text",
                    "content": f"Error generating auditory content: {str(e)}\n\nOriginal text (partial):\n\n{document_text[:500]}...",
                    "caption": "Error Processing Document"
                }
            ]
        }

@app.route('/api/files/<file_id>/process-auditory', methods=['OPTIONS', 'POST'])
def process_auditory(file_id):
    """Process a document for auditory learners by generating spoken-friendly content and audio file."""
    if request.method == 'OPTIONS':
        return create_preflight_response()
        
    try:
        print(f"Processing auditory content for file ID: {file_id}")
        
        # Check if OpenAI API key is set
        if not openai_api_key:
            print("OpenAI API key is not set")
            return jsonify({"error": "OpenAI API key is not configured"}), 500
            
        # Check if Firebase is initialized
        if not db or not bucket:
            print("Firebase is not properly initialized")
            return jsonify({"error": "Firebase is not properly initialized"}), 500
        
        # Get the document from Firestore
        doc_ref = db.collection("files").document(file_id)
        doc = doc_ref.get()
        
        if not doc.exists:
            print(f"Document with ID {file_id} not found")
            return jsonify({"error": "Document not found"}), 404
            
        doc_data = doc.to_dict()
        storage_path = doc_data.get("storagePath")
        
        if not storage_path:
            print(f"Storage path not found for document {file_id}")
            return jsonify({"error": "Storage path not found"}), 400
            
        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            
        try:
            # Download file
            temp_path = os.path.join(temp_dir, f"{file_id}_{os.path.basename(storage_path)}")
            blob = bucket.blob(storage_path)
            blob.download_to_filename(temp_path)
            
            # Extract text from document
            document_text = extract_text_from_document(temp_path)
            
            if not document_text:
                raise ValueError("Failed to extract text from document")
                
            # Generate auditory content and audio file
            content = generate_auditory_content(document_text)
            
            if not content:
                raise ValueError("Failed to generate auditory content")
            
            # Store the processed content in Firestore
            doc_ref.update({
                "auditoryContent": content,
                "auditoryGeneratedAt": firestore.SERVER_TIMESTAMP
            })
            
            return jsonify({
                "success": True,
                "content": content,
                "audioUrl": content.get("audioUrl")
            })
            
        except Exception as inner_error:
            print(f"Error during file processing: {str(inner_error)}")
            raise inner_error
            
        finally:
            # Clean up temp file
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception as cleanup_error:
                    print(f"Error cleaning up temp file: {str(cleanup_error)}")
        
    except Exception as e:
        print(f"Error processing auditory content: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

def generate_kinesthetic_content(document_text):
    """
    Generate interactive, hands-on learning activities for kinesthetic learners using OpenAI
    """
    if not openai_api_key:
        print("Using mock response for kinesthetic content (No API key)")
        return {
            "title": "Interactive Learning Activities",
            "description": "Learn through hands-on activities and physical engagement.",
            "activities": [
                {
                    "title": "Sample Activity",
                    "description": "A hands-on exercise to understand the concept.",
                    "materials": ["Item 1", "Item 2"],
                    "steps": ["Step 1", "Step 2", "Step 3"],
                    "tips": ["Tip 1", "Tip 2"],
                    "reflection": ["Question 1", "Question 2"]
                }
            ]
        }
    
    try:
        # Truncate text if too long (OpenAI has token limits)
        max_chars = 14000
        truncated_text = document_text[:max_chars] if len(document_text) > max_chars else document_text
        
        # Call OpenAI API to generate kinesthetic activities
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": """You are an AI tutor specialized in creating hands-on, interactive learning activities for kinesthetic learners. 
                Transform educational content into engaging physical exercises, experiments, and real-world applications.
                
                For each activity, you MUST provide:
                1. Title: A clear, action-oriented title
                2. Description: Brief overview of what will be learned
                3. Materials Needed: List of required items (at least 2)
                4. Steps: Detailed numbered instructions (at least 3 steps)
                5. Tips: Practical advice for success (at least 2 tips)
                6. Reflection Questions: Thought-provoking questions (at least 2 questions)
                
                Format each activity section with these exact headings and format:
                Title: [Activity Title]
                Description: [Brief overview]
                Materials:
                - [First item]
                - [Second item]
                Steps:
                1. [First step]
                2. [Second step]
                3. [Third step]
                Tips:
                - [First tip]
                - [Second tip]
                Reflection Questions:
                - [First question]
                - [Second question]
                
                Create 2-3 activities that are:
                - Hands-on and physically engaging
                - Using readily available materials
                - Safe and appropriate
                - Connected to real-world applications
                - Suitable for both individual and group work"""},
                {"role": "user", "content": f"Create interactive, hands-on learning activities for this content:\n\n{truncated_text}"}
            ],
            max_tokens=2000,
            temperature=0.7
        )
        
        # Extract the formatted text
        activities_text = response.choices[0].message.content
        
        # Parse the activities into structured format
        activities = []
        current_activity = None
        current_section = None
        
        lines = activities_text.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            if not line:
                i += 1
                continue
                
            # Check for new activity
            if line.startswith('Title:'):
                if current_activity and current_activity['title']:  # Save previous activity if exists
                    activities.append(current_activity)
                current_activity = {
                    'title': line.split(':', 1)[1].strip(),
                    'description': '',
                    'materials': [],
                    'steps': [],
                    'tips': [],
                    'reflection': []
                }
                current_section = 'title'
            
            # Check for sections
            elif line.startswith('Description:'):
                current_section = 'description'
                current_activity['description'] = line.split(':', 1)[1].strip()
            elif line.startswith('Materials:'):
                current_section = 'materials'
            elif line.startswith('Steps:'):
                current_section = 'steps'
            elif line.startswith('Tips:'):
                current_section = 'tips'
            elif line.startswith('Reflection Questions:'):
                current_section = 'reflection'
            
            # Process content based on current section
            elif current_activity and current_section:
                content = line
                
                # Handle numbered steps
                if line[0].isdigit() and '. ' in line:
                    content = line.split('. ', 1)[1].strip()
                # Handle bullet points
                elif line.startswith('- ') or line.startswith('* '):
                    content = line[2:].strip()
                # Handle continuation of description
                elif current_section == 'description' and not line.startswith(('Materials:', 'Steps:', 'Tips:', 'Reflection Questions:', 'Title:')):
                    current_activity['description'] += ' ' + content
                    i += 1
                    continue
                
                # Skip if content is empty or just a bullet/number
                if not content or content in ['-', '*'] or content.isdigit():
                    i += 1
                    continue
                
                # Add content to appropriate section
                if current_section == 'materials':
                    current_activity['materials'].append(content)
                elif current_section == 'steps':
                    current_activity['steps'].append(content)
                elif current_section == 'tips':
                    current_activity['tips'].append(content)
                elif current_section == 'reflection':
                    current_activity['reflection'].append(content)
            
            i += 1
        
        # Add the last activity if exists
        if current_activity and current_activity['title']:
            activities.append(current_activity)
        
        # Validate activities
        validated_activities = []
        for activity in activities:
            if (activity['title'] and 
                activity['description'] and 
                len(activity['materials']) >= 1 and 
                len(activity['steps']) >= 2 and 
                len(activity['tips']) >= 1 and 
                len(activity['reflection']) >= 1):
                validated_activities.append(activity)
        
        if not validated_activities:
            print("No valid activities were generated")
            return {
                "title": "Interactive Learning Activities",
                "description": "Learn through hands-on activities and physical engagement.",
                "activities": [{
                    "title": "Activity Generation Error",
                    "description": "Please try regenerating the activities.",
                    "materials": ["Paper", "Pen"],
                    "steps": ["Review the content carefully", "Take notes on key concepts", "Practice explaining the concepts"],
                    "tips": ["Focus on understanding one concept at a time", "Try to relate concepts to real-world examples"],
                    "reflection": ["What are the main ideas you learned?", "How can you apply this knowledge?"]
                }]
            }
        
        return {
            "title": "Interactive Learning Activities",
            "description": "Learn through hands-on activities and physical engagement.",
            "activities": validated_activities
        }
        
    except Exception as e:
        print(f"Error generating kinesthetic content: {e}")
        return {
            "title": "Interactive Learning Activities (Error)",
            "description": "There was an error processing this document.",
            "activities": [{
                "title": "Temporary Activity",
                "description": "While we fix the error, here's a general learning activity.",
                "materials": ["Paper", "Pen", "Textbook or learning materials"],
                "steps": [
                    "Review the main concepts in your materials",
                    "Create a mind map or diagram of key ideas",
                    "Practice explaining the concepts out loud"
                ],
                "tips": [
                    "Take breaks between concepts",
                    "Try to relate ideas to real-world examples"
                ],
                "reflection": [
                    "What was the most challenging concept to understand?",
                    "How can you apply what you learned?"
                ]
            }]
        }

@app.route('/api/files/<file_id>/process-kinesthetic', methods=['OPTIONS', 'POST'])
def process_kinesthetic(file_id):
    """Process a document for kinesthetic learners by generating interactive activities."""
    if request.method == 'OPTIONS':
        return create_preflight_response()
        
    try:
        print(f"Processing kinesthetic content for file ID: {file_id}")
        
        # Get the document from Firestore
        doc_ref = db.collection("files").document(file_id)
        doc = doc_ref.get()
        
        if not doc.exists:
            print(f"Document with ID {file_id} not found")
            return jsonify({"error": "Document not found"}), 404
            
        doc_data = doc.to_dict()
        storage_path = doc_data.get("storagePath")
        
        if not storage_path:
            print(f"Storage path not found for document {file_id}")
            return jsonify({"error": "Storage path not found"}), 400
            
        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            
        # Download file
        temp_path = os.path.join(temp_dir, f"{file_id}_{os.path.basename(storage_path)}")
        blob = bucket.blob(storage_path)
        blob.download_to_filename(temp_path)
        
        # Extract text from document
        document_text = extract_text_from_document(temp_path)
        
        # Clean up temp file
        if os.path.exists(temp_path):
            os.remove(temp_path)
            
        # Generate kinesthetic content
        content = generate_kinesthetic_content(document_text)
        
        # Store the processed content in Firestore
        doc_ref.update({
            "kinestheticContent": content,
            "kinestheticGeneratedAt": firestore.SERVER_TIMESTAMP
        })
        
        return jsonify({
            "success": True,
            "content": content
        })
        
    except Exception as e:
        print(f"Error processing kinesthetic content: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

def generate_visual_content(document_text):
    """Generate visual learning suggestions and concept explanations for visual learning."""
    if not openai_api_key:
        print("Using mock response for visual content (No API key)")
        return {
            "success": True,
            "content": {
                "title": "Visual Learning Materials",
                "description": "Learn through diagrams, concept maps, and visual representations.",
                "suggestions": [
                    "Create mind maps connecting key concepts from the document",
                    "Use color-coding to highlight related information",
                    "Draw timelines to visualize sequences and processes",
                    "Convert text information into diagrams, charts, or graphs"
                ],
                "explanations": [
                    {
                        "title": "Concept Mapping",
                        "text": "Organize information visually by connecting related ideas with lines or arrows",
                        "image": "/static/images/placeholder.png"
                    },
                    {
                        "title": "Visual Hierarchies",
                        "text": "Represent information in a top-down structure showing relationships between main topics and subtopics",
                        "image": "/static/images/placeholder.png"
                    },
                    {
                        "title": "Color Coding",
                        "text": "Use colors systematically to categorize information and highlight patterns",
                        "image": "/static/images/placeholder.png"
                    }
                ]
            }
        }
    
    try:
        # Truncate text if too long
        max_chars = 14000
        truncated_text = document_text[:max_chars] if len(document_text) > max_chars else document_text
        
        # Generate visual learning suggestions
        suggestions_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": """Create visual learning suggestions for the provided content.
                Generate a list of 4-6 specific ways a visual learner could engage with this content.
                Focus on techniques like mind mapping, diagramming, sketching, and visual note-taking."""},
                {"role": "user", "content": f"Create visual learning suggestions for this content:\n\n{truncated_text}"}
            ],
            temperature=0.7
        )
        
        # Extract suggestions as a list
        suggestions_content = suggestions_response.choices[0].message.content.strip()
        suggestions = []
        
        for line in suggestions_content.split('\n'):
            # Remove numbered bullets or dashes if present
            clean_line = re.sub(r'^\d+\.\s*|^-\s*', '', line.strip())
            if clean_line:
                suggestions.append(clean_line)
        
        # Extract main concepts from text
        concepts_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": """Extract 3-5 main concepts or topics from the provided text.
                For each concept:
                1. Provide a short title (1-3 words)
                2. Write a brief explanation (50-70 words)
                
                Return a list of concepts in JSON format with:
                - 'title': The concept name
                - 'text': A brief explanation
                """},
                {"role": "user", "content": f"Extract key concepts from this content:\n\n{truncated_text}"}
            ],
            response_format={ "type": "json_object" },
            temperature=0.7
        )
        
        # Parse the response to get explanations
        try:
            concepts_data = json.loads(concepts_response.choices[0].message.content)
            explanations = concepts_data.get("concepts", [])
            
            # If response doesn't have the expected format, try parsing differently
            if not explanations and isinstance(concepts_data, list):
                explanations = concepts_data
                
            # Process explanations to ensure correct format
            processed_explanations = []
            for exp in explanations:
                if isinstance(exp, dict) and 'title' in exp:
                    # Set a default image path
                    if 'image' not in exp:
                        exp['image'] = '/static/images/placeholder.png'
                    processed_explanations.append(exp)
            
            if processed_explanations:
                return {
                    "success": True,
                    "content": {
                        "title": "Visual Learning Materials",
                        "description": "Learn through diagrams, concept maps, and visual representations.",
                        "suggestions": suggestions,
                        "explanations": processed_explanations
                    }
                }
            else:
                raise ValueError("Failed to parse explanations from OpenAI response")
                
        except (json.JSONDecodeError, KeyError, ValueError) as e:
            print(f"Error parsing concepts response: {e}")
            print(f"Raw response: {concepts_response.choices[0].message.content}")
            raise
            
    except Exception as e:
        print(f"Error generating visual content: {e}")
        return {
            "success": False,
            "error": str(e)
        }

@app.route('/api/files/<file_id>/process-visual', methods=['OPTIONS', 'POST'])
def process_visual(file_id):
    """Process a document for visual learners by generating visual learning suggestions."""
    if request.method == 'OPTIONS':
        return create_preflight_response()
        
    try:
        print(f"Processing visual content for file ID: {file_id}")
        
        # Get the document from Firestore
        doc_ref = db.collection("files").document(file_id)
        doc = doc_ref.get()
        
        if not doc.exists:
            print(f"Document with ID {file_id} not found")
            return jsonify({"error": "Document not found"}), 404
            
        doc_data = doc.to_dict()
        storage_path = doc_data.get("storagePath")
        
        if not storage_path:
            print(f"Storage path not found for document {file_id}")
            return jsonify({"error": "Storage path not found"}), 400
            
        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            
        # Download file
        temp_path = os.path.join(temp_dir, f"{file_id}_{os.path.basename(storage_path)}")
        blob = bucket.blob(storage_path)
        blob.download_to_filename(temp_path)
        
        # Extract text from document
        document_text = extract_text_from_document(temp_path)
        
        # Clean up temp file
        if os.path.exists(temp_path):
            os.remove(temp_path)
            
        # Generate visual content (now just suggestions)
        content = generate_visual_content(document_text)
        
        # Validate explanations format
        if "explanations" in content:
            # Ensure explanations is an array
            if not isinstance(content["explanations"], list):
                content["explanations"] = []
            
            # Validate each explanation
            valid_explanations = []
            for expl in content["explanations"]:
                if isinstance(expl, dict) and "title" in expl and "image" in expl and "text" in expl:
                    valid_explanations.append({
                        "title": str(expl["title"]),
                        "image": str(expl["image"]),
                        "text": str(expl["text"])
                    })
            content["explanations"] = valid_explanations
        
        # Validate suggestions format
        if "suggestions" in content:
            # Ensure suggestions is an array of strings
            if not isinstance(content["suggestions"], list):
                content["suggestions"] = []
            
            # Convert non-string suggestions to strings
            content["suggestions"] = [str(s) for s in content["suggestions"]]
        
        # Store the processed content in Firestore
        doc_ref.update({
            "visualContent": content,
            "visualGeneratedAt": firestore.SERVER_TIMESTAMP
        })
        
        return jsonify({
            "success": True,
            "content": content
        })
        
    except Exception as e:
        print(f"Error processing visual content: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

def generate_concept_svg(concept_name):
    """Generate a simple SVG icon based on the concept name."""
    # Generate a consistent color based on the concept name
    hash_value = int(hashlib.md5(concept_name.encode()).hexdigest(), 16)
    
    # Generate colors for the concept
    hue = hash_value % 360
    primary_color = f"hsl({hue}, 70%, 60%)"  # Primary color
    secondary_color = f"hsl({(hue + 40) % 360}, 60%, 70%)"  # Secondary color
    accent_color = f"hsl({(hue + 180) % 360}, 70%, 50%)"  # Accent color
    
    # Get the first letter of each word in the concept (max 2)
    words = concept_name.split()
    initial = words[0][0].upper()
    if len(words) > 1:
        initial += words[1][0].upper()
    else:
        initial = initial[:2]  # Take at most 2 chars
    
    # Create a unique filename
    safe_name = re.sub(r'[^a-z0-9]', '-', concept_name.lower())
    filename = f"{safe_name}.svg"
    
    # Create the SVG file path
    svg_path = os.path.join('frontend', 'public', 'images', filename)
    
    # Generate a simple SVG
    svg_content = f"""<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100" viewBox="0 0 100 100">
  <rect width="100" height="100" fill="#f5f9ff" rx="10" ry="10"/>
  <circle cx="50" cy="50" r="30" fill="{primary_color}"/>
  <path d="M30 70 L70 30 M30 30 L70 70" stroke="{secondary_color}" stroke-width="5" stroke-linecap="round"/>
  <text x="50" y="55" font-family="Arial" font-size="20" text-anchor="middle" fill="white" font-weight="bold">{initial}</text>
</svg>"""
    
    # Ensure the directory exists
    os.makedirs(os.path.dirname(svg_path), exist_ok=True)
    
    # Write the SVG file
    with open(svg_path, 'w') as f:
        f.write(svg_content)
    
    return f"/images/{filename}"

@app.route('/api/generate-visuals', methods=['POST'])
def generate_visuals():
    """
    Extract visual concepts from submitted text.
    Accepts JSON with a 'text' field containing the document content.
    Returns a list of concept objects with title, text, and image path.
    """
    if request.method == 'OPTIONS':
        return create_preflight_response()
        
    try:
        # Get text from the request
        data = request.get_json()
        
        if not data or 'text' not in data:
            return jsonify({"error": "No text provided"}), 400
            
        # Extract text from the request
        document_text = data['text']
        
        if not document_text or len(document_text.strip()) < 50:
            return jsonify({"error": "Text is too short for analysis"}), 400
            
        # Generate visual concepts
        if not openai_api_key:
            print("Using mock response for visual concepts (No API key)")
            
            # Return sample data
            return jsonify({
                "success": True,
                "explanations": [
                    {
                        "title": "Visual Learning",
                        "image": "/images/visual-learning.svg",
                        "text": "Learning through seeing and visualizing concepts."
                    },
                    {
                        "title": "Mind Maps",
                        "image": "/images/mind-map.svg",
                        "text": "Organizing information in a visual hierarchy to see relationships."
                    },
                    {
                        "title": "Color Coding",
                        "image": "/images/color-coding.svg",
                        "text": "Using different colors to categorize and remember information."
                    }
                ]
            })
        
        # Truncate text if too long (OpenAI has token limits)
        max_chars = 14000
        truncated_text = document_text[:max_chars] if len(document_text) > max_chars else document_text
        
        # Extract main concepts from text
        concepts_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": """Extract 4-6 main concepts or topics from the provided text.
                For each concept:
                1. Provide a short title (1-3 words)
                2. Write a brief 1-2 sentence explanation of the concept
                3. Suggest an appropriate image filename in the format "/images/concept-name.svg"
                
                Return a JSON object with an array called 'explanations', where each item has:
                - 'title': The concept name
                - 'text': A brief explanation (max 120 characters)
                - 'image': The image path
                
                Keep explanations clear, concise, and suitable for visual learning."""},
                {"role": "user", "content": f"Extract key concepts from this content:\n\n{truncated_text}"}
            ],
            response_format={ "type": "json_object" },
            temperature=0.7
        )
        
        # Parse response
        concepts_data = json.loads(concepts_response.choices[0].message.content)
        
        # Map concept images to existing SVGs where possible
        standard_images = {
            "visual learning": "/images/visual-learning.svg",
            "mind map": "/images/mind-map.svg", 
            "color coding": "/images/color-coding.svg"
        }
        
        # Process explanations to use standard images where appropriate
        explanations = concepts_data.get("explanations", [])
        for explanation in explanations:
            title_lower = explanation["title"].lower()
            
            # Check for standard images first
            image_matched = False
            for key, image_path in standard_images.items():
                if key in title_lower:
                    explanation["image"] = image_path
                    image_matched = True
                    break
            
            # Generate a custom SVG if no standard image matches
            if not image_matched:
                explanation["image"] = generate_concept_svg(explanation["title"])
        
        return jsonify({
            "success": True,
            "explanations": explanations
        })
        
    except Exception as e:
        print(f"Error generating visual concepts: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/api/visual-concepts', methods=['GET'])
def get_visual_concepts():
    """
    Return visual concepts for a document
    Accepts a fileId query parameter
    """
    try:
        file_id = request.args.get('fileId')
        if not file_id:
            return jsonify({"error": "No fileId provided"}), 400
            
        # Get the document from Firestore
        doc_ref = db.collection("files").document(file_id)
        doc = doc_ref.get()
        
        if not doc.exists:
            return jsonify({"error": "Document not found"}), 404
            
        doc_data = doc.to_dict()
        storage_path = doc_data.get("storagePath")
        
        if not storage_path:
            return jsonify({"error": "Storage path not found"}), 400
            
        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            
        # Download file
        temp_path = os.path.join(temp_dir, f"{file_id}_{os.path.basename(storage_path)}")
        blob = bucket.blob(storage_path)
        blob.download_to_filename(temp_path)
        
        # Extract text from document
        document_text = extract_text_from_document(temp_path)
        
        # Clean up temp file
        if os.path.exists(temp_path):
            os.remove(temp_path)
            
        # Generate sample concepts (in production, use a real NLP model)
        if not openai_api_key:
            # Sample data for testing - use placeholder images
            return jsonify([
                {
                    "title": "Convolutional Neural Network",
                    "description": "Utilized for classification and computer vision tasks.",
                    "image": "/static/images/placeholder.png"
                },
                {
                    "title": "Recurrent Neural Network",
                    "description": "For sequential data processing, like text or time series.",
                    "image": "/static/images/placeholder.png"
                },
                {
                    "title": "Transfer Learning",
                    "description": "Using pre-trained models and fine-tuning for new tasks.",
                    "image": "/static/images/placeholder.png"
                }
            ])
        
        # Use OpenAI to generate concepts
        try:
            # Truncate text if too long (OpenAI has token limits)
            max_chars = 14000
            truncated_text = document_text[:max_chars] if len(document_text) > max_chars else document_text
            
            # Extract main concepts from text
            concepts_response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": """Extract 4-6 main concepts or topics from the provided text.
                    For each concept:
                    1. Provide a short title (1-3 words)
                    2. Write a brief 1-2 sentence explanation of the concept
                    
                    Return a JSON array where each item has:
                    - 'title': The concept name
                    - 'description': A brief explanation (max 120 characters)
                    """},
                    {"role": "user", "content": f"Extract key concepts from this content:\n\n{truncated_text}"}
                ],
                response_format={ "type": "json_object" },
                temperature=0.7
            )
            
            # Parse response
            data = json.loads(concepts_response.choices[0].message.content)
            concepts = data if isinstance(data, list) else data.get("concepts", [])
            
            # Get images for each concept using SerpAPI
            for concept in concepts:
                try:
                    # Check for SerpAPI key
                    if not serpapi_key:
                        concept["image"] = "/static/images/placeholder.png"
                        continue
                        
                    if GoogleSearch is None:
                        concept["image"] = "/static/images/placeholder.png"
                        continue
                        
                    # Use SerpAPI to get image
                    search = GoogleSearch({
                        "q": f"{concept['title']} diagram educational",
                        "tbm": "isch",  # image search
                        "num": 1,
                        "safe": "active",  # safe search
                        "api_key": serpapi_key
                    })
                    results = search.get_dict()
                    
                    if results.get("images_results"):
                        concept["image"] = results["images_results"][0]["original"]
                    else:
                        concept["image"] = "/static/images/placeholder.png"
                        
                except Exception as img_error:
                    print(f"Error getting image for {concept['title']}: {img_error}")
                    concept["image"] = "/static/images/placeholder.png"
                
            return jsonify(concepts)
            
        except Exception as e:
            print(f"Error generating concepts with OpenAI: {e}")
            # Fallback to sample data with placeholder images
            return jsonify([
                {
                    "title": "Convolutional Neural Network",
                    "description": "Utilized for classification and computer vision tasks.",
                    "image": "/static/images/placeholder.png"
                },
                {
                    "title": "Recurrent Neural Network", 
                    "description": "For sequential data processing, like text or time series.",
                    "image": "/static/images/placeholder.png"
                },
                {
                    "title": "Transfer Learning",
                    "description": "Using pre-trained models and fine-tuning for new tasks.",
                    "image": "/static/images/placeholder.png"
                }
            ])
        
    except Exception as e:
        print(f"Error in visual concepts endpoint: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/api/generate-image', methods=['POST'])
def generate_image():
    """Redirects to SerpAPI-based image search instead of using DALL-E"""
    try:
        data = request.json
        description = data.get('description')
        
        if not description:
            return jsonify({"error": "Description is required"}), 400
            
        # Redirect to SerpAPI search
        return image_for_topic()
        
    except Exception as e:
        print(f"Error redirecting image generation: {str(e)}")
        return jsonify({
            "success": False,
            "error": str(e)
        }), 500

@app.route('/api/image-for-topic', methods=['POST'])
def image_for_topic():
    """
    Search Google Images for diagrams related to a topic using SerpAPI
    """
    topic = request.json.get('topic', '').strip()

    if not topic:
        return jsonify({"error": "No topic provided"}), 400
        
    # Get the SerpAPI key from environment variables
    serpapi_key = os.getenv("SERPAPI_API_KEY")
    
    if not serpapi_key:
        return jsonify({"error": "SerpAPI key not configured"}), 500
        
    if GoogleSearch is None:
        return jsonify({"error": "SerpAPI package not installed. Please run: pip install google-search-results"}), 500

    try:
        search = GoogleSearch({
            "q": f"{topic} diagram educational",
            "tbm": "isch",  # image search
            "num": 1,
            "safe": "active",  # safe search
            "api_key": serpapi_key
        })
        results = search.get_dict()
        
        if not results.get("images_results"):
            return jsonify({"error": "No images found", "image": "/static/images/placeholder.png"}), 404
            
        image_url = results["images_results"][0]["original"]
        return jsonify({
            "success": true,
            "image_url": image_url
        })
        
    except Exception as e:
        print(f"Error searching for image: {str(e)}")
        return jsonify({
            "success": false,
            "error": str(e),
            "image_url": "/static/images/placeholder.png"
        }), 500

@app.route('/api/files/<file_id>/process-reading-writing-consistent', methods=['OPTIONS', 'POST'])
def process_reading_writing_consistent(file_id):
    """Process a document for reading/writing learners with consistent output across web and files."""
    # Handle preflight request
    if request.method == 'OPTIONS':
        return create_preflight_response()
        
    try:
        print(f"Processing reading/writing content (CONSISTENT) for file ID: {file_id}")
        
        # Get the document from Firestore
        doc_ref = db.collection("files").document(file_id)
        doc = doc_ref.get()
        
        if not doc.exists:
            print(f"Document with ID {file_id} not found")
            return jsonify({"error": "Document not found"}), 404
            
        doc_data = doc.to_dict()
        storage_path = doc_data.get("storagePath")
        
        if not storage_path:
            print(f"Storage path not found for document {file_id}")
            return jsonify({"error": "Storage path not found"}), 400
            
        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(os.getcwd(), 'temp_processing')
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
            
        # Download file
        temp_path = os.path.join(temp_dir, f"{file_id}_{os.path.basename(storage_path)}")
        blob = bucket.blob(storage_path)
        blob.download_to_filename(temp_path)
        
        # Extract text from document
        document_text = extract_text_from_document(temp_path)
        
        # Clean up temp file
        if os.path.exists(temp_path):
            os.remove(temp_path)
            
        # Generate reading/writing optimized content
        print("Generating consistent content for web and files...")
        content = generate_reading_writing_content(document_text)
        
        # Get the content that will be displayed on the web - EXACTLY the same will be used for files
        web_content = content["elements"][0]["content"]
        print(f"Content first 100 chars: {web_content[:100]}")
        
        # Create DOCX and PDF with EXACTLY the same content
        docx_file = create_docx_study_guide(web_content)
        pdf_file = create_pdf_study_guide(web_content)
        
        # Upload processed documents to Firebase Storage
        docx_filename = f"reading_writing_{file_id}.docx"
        pdf_filename = f"reading_writing_{file_id}.pdf"
        
        docx_storage_path = f"processed/{doc_data.get('userId', 'unknown')}/{docx_filename}"
        pdf_storage_path = f"processed/{doc_data.get('userId', 'unknown')}/{pdf_filename}"
        
        # Upload DOCX
        docx_blob = bucket.blob(docx_storage_path)
        docx_blob.upload_from_file(docx_file, content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        docx_blob.make_public()
        
        # Upload PDF
        pdf_blob = bucket.blob(pdf_storage_path)
        pdf_blob.upload_from_file(pdf_file, content_type='application/pdf')
        pdf_blob.make_public()
        
        # Store the processed content in Firestore
        doc_ref.update({
            "readingWritingContent": content,
            "readingWritingDocxUrl": docx_blob.public_url,
            "readingWritingPdfUrl": pdf_blob.public_url,
            "readingWritingGeneratedAt": firestore.SERVER_TIMESTAMP
        })
        
        return jsonify({
            "success": True,
            "content": content,
            "docxUrl": docx_blob.public_url,
            "pdfUrl": pdf_blob.public_url
        })
        
    except Exception as e:
        print(f"Error processing reading/writing content: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/api/user/login-streak', methods=['POST'])
def update_login_streak():
    """Update user's login streak in Firestore"""
    try:
        data = request.json
        if not data:
            return jsonify({"error": "Invalid request: No JSON data provided"}), 400
            
        email = data.get("email")
        
        if not email:
            return jsonify({"error": "Email is required"}), 400
        
        # Find user in Firestore
        users_ref = db.collection("users")
        user_query = users_ref.where("email", "==", email).limit(1)
        user_docs = list(user_query.stream())
        
        if not user_docs:
            return jsonify({"error": "User not found"}), 404
        
        user_id = user_docs[0].id
        user_data = user_docs[0].to_dict()
        
        # Get current time in the user's timezone (default to UTC)
        # In a production app, you might want to store user's timezone
        current_time = datetime.now()
        
        # Get last login date if exists
        last_login_date = None
        current_streak = 0
        longest_streak = 0
        
        if "lastLoginDate" in user_data:
            last_login_str = user_data.get("lastLoginDate")
            try:
                last_login_date = datetime.fromisoformat(last_login_str)
            except (ValueError, TypeError):
                # If date format is invalid, treat as if no previous login
                pass
                
        if "currentStreak" in user_data:
            current_streak = user_data.get("currentStreak", 0)
            
        if "longestStreak" in user_data:
            longest_streak = user_data.get("longestStreak", 0)
        
        # Calculate if streak should increase, reset, or stay the same
        streak_data = {}
        
        # First login ever
        if last_login_date is None:
            streak_data = {
                "lastLoginDate": current_time.isoformat(),
                "currentStreak": 1,
                "longestStreak": 1
            }
        else:
            # Calculate days between logins
            days_difference = (current_time.date() - last_login_date.date()).days
            
            if days_difference == 0:
                # Already logged in today, don't update streak
                streak_data = {
                    "lastLoginDate": current_time.isoformat()
                }
            elif days_difference == 1:
                # Consecutive day login, increase streak
                new_streak = current_streak + 1
                streak_data = {
                    "lastLoginDate": current_time.isoformat(),
                    "currentStreak": new_streak,
                    "longestStreak": max(new_streak, longest_streak)
                }
            else:
                # Streak broken, reset to 1
                streak_data = {
                    "lastLoginDate": current_time.isoformat(),
                    "currentStreak": 1,
                    "longestStreak": max(1, longest_streak)
                }
        
        # Update user record in Firestore
        db.collection("users").document(user_id).update(streak_data)
        
        # Return updated streak information
        return jsonify({
            "success": True,
            "streak": streak_data.get("currentStreak", current_streak),
            "longestStreak": streak_data.get("longestStreak", longest_streak)
        })
    except Exception as e:
        print(f"Error updating login streak: {str(e)}")
        return jsonify({"error": f"Failed to update login streak: {str(e)}"}), 500

@app.route('/api/user/record-time', methods=['POST'])
def record_session_time():
    """Record user's session time and update total time spent"""
    try:
        data = request.json
        if not data:
            return jsonify({"error": "Invalid request: No JSON data provided"}), 400
            
        email = data.get("email")
        duration_minutes = data.get("durationMinutes")
        session_start = data.get("sessionStart")
        session_end = data.get("sessionEnd")
        activity_type = data.get("activityType", "session")
        
        if not email:
            return jsonify({"error": "Email is required"}), 400
        
        if duration_minutes is None:
            return jsonify({"error": "Duration is required"}), 400
        
        # Find user in Firestore
        users_ref = db.collection("users")
        user_query = users_ref.where("email", "==", email).limit(1)
        user_docs = list(user_query.stream())
        
        if not user_docs:
            return jsonify({"error": "User not found"}), 404
        
        user_id = user_docs[0].id
        user_data = user_docs[0].to_dict()
        
        # Initialize stats if not present
        if "stats" not in user_data:
            user_data["stats"] = {}
        
        if "totalTimeSpent" not in user_data["stats"]:
            user_data["stats"]["totalTimeSpent"] = 0
            
        if "sessions" not in user_data["stats"]:
            user_data["stats"]["sessions"] = []
            
        if "weeklyStats" not in user_data["stats"]:
            user_data["stats"]["weeklyStats"] = {}
            
        # Update total time spent
        current_total = user_data["stats"]["totalTimeSpent"]
        new_total = current_total + duration_minutes
        
        # Create session record
        session_record = {
            "startTime": session_start,
            "endTime": session_end,
            "durationMinutes": duration_minutes,
            "type": activity_type
        }
        
        # Add new session to history (limit to last 50 sessions)
        sessions = user_data["stats"]["sessions"]
        sessions.append(session_record)
        if len(sessions) > 50:
            sessions = sessions[-50:]  # Keep only the last 50 sessions
            
        # Update weekly stats
        weekly_stats = user_data["stats"].get("weeklyStats", {})
        current_week_data = None
        
        try:
            session_datetime = datetime.fromisoformat(session_start)
            # Get ISO year and week number
            year, week_num, _ = session_datetime.isocalendar()
            week_key = f"{year}-W{week_num:02d}"
            
            # Check if lastWeekTracked exists and if we're in a new week - for automatic reset
            last_week_tracked = user_data["stats"].get("lastWeekTracked", None)
            if last_week_tracked and last_week_tracked != week_key:
                # We've moved to a new week, but we keep the old weeks' data in the weeklyStats object
                # The frontend will only show current week
                print(f"User moved to a new week: {last_week_tracked} -> {week_key}")
            
            # Update lastWeekTracked
            user_data["stats"]["lastWeekTracked"] = week_key
            
            # Ensure the current week exists
            if week_key not in weekly_stats:
                weekly_stats[week_key] = {
                    "totalMinutes": 0,
                    "activeDays": [],
                    "activities": {}
                }
                
            weekly_stats[week_key]["totalMinutes"] += duration_minutes
            current_week_data = weekly_stats[week_key]
            
            # Add day to active days if not already there
            day_str = session_datetime.strftime("%Y-%m-%d")
            if day_str not in weekly_stats[week_key]["activeDays"]:
                weekly_stats[week_key]["activeDays"].append(day_str)
                
            # Update activity breakdown
            if activity_type not in weekly_stats[week_key]["activities"]:
                weekly_stats[week_key]["activities"][activity_type] = 0
            weekly_stats[week_key]["activities"][activity_type] += duration_minutes
            
        except (ValueError, TypeError):
            print(f"Error parsing date: {session_start}")
        
        # Update user record in Firestore
        db.collection("users").document(user_id).update({
            "stats.totalTimeSpent": new_total,
            "stats.sessions": sessions,
            "stats.weeklyStats": weekly_stats,
            "stats.lastWeekTracked": user_data["stats"].get("lastWeekTracked"),
            "stats.currentWeekTime": current_week_data["totalMinutes"] if current_week_data else 0,
            "lastActive": datetime.now().isoformat()
        })
        
        # Trigger leaderboard update if time is significant (every 10 minutes)
        # This prevents updating the leaderboard too frequently while still keeping it updated
        if duration_minutes >= 10:
            try:
                # Use a background thread to avoid slowing down the response
                from threading import Thread
                Thread(target=lambda: requests.post(
                    f"{request.url_root}api/leaderboard/update", 
                    json={}, 
                    timeout=10
                )).start()
            except Exception as e:
                print(f"Error triggering leaderboard update: {str(e)}")
                # Don't fail the main request if this fails
        
        return jsonify({
            "success": True,
            "totalTimeSpent": new_total,
            "sessionDuration": duration_minutes,
            "weeklyStats": weekly_stats,
            "currentWeekTime": current_week_data["totalMinutes"] if current_week_data else 0
        })
    except Exception as e:
        print(f"Error recording session time: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Failed to record session time: {str(e)}"}), 500

@app.route('/api/leaderboard', methods=['GET'])
def get_leaderboard():
    """Get leaderboard data for top streaks and learning time"""
    try:
        # Get leaderboard document from Firestore
        leaderboard_ref = db.collection("leaderboard").document("stats")
        leaderboard_doc = leaderboard_ref.get()
        
        if not leaderboard_doc.exists:
            # Initialize leaderboard if it doesn't exist
            leaderboard_data = {
                "topStreaks": [],
                "topLearningTime": []
            }
            db.collection("leaderboard").document("stats").set(leaderboard_data)
        else:
            leaderboard_data = leaderboard_doc.to_dict()
            
        return jsonify({
            "success": True,
            "leaderboard": leaderboard_data
        })
    except Exception as e:
        print(f"Error getting leaderboard: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Failed to get leaderboard: {str(e)}"}), 500

@app.route('/api/leaderboard/update', methods=['POST'])
def update_leaderboard():
    """Update leaderboard with all users' stats"""
    try:
        # Get all users
        users_ref = db.collection("users")
        users = list(users_ref.stream())
        
        # Prepare streak and time data
        streak_data = []
        time_data = []
        
        for user_doc in users:
            user_data = user_doc.to_dict()
            user_id = user_doc.id
            
            # Skip users without necessary data
            if not user_data.get("email"):
                continue
                
            # Get display name or use email prefix
            display_name = user_data.get("displayName") or user_data.get("name")
            if not display_name and user_data.get("email"):
                display_name = user_data["email"].split('@')[0]
            
            # Get streak data
            current_streak = user_data.get("currentStreak", 0)
            if current_streak > 0:
                streak_data.append({
                    "userId": user_id,
                    "username": display_name,
                    "streak": current_streak
                })
            
            # Get time spent data
            stats = user_data.get("stats", {})
            total_time = stats.get("totalTimeSpent", 0)
            if total_time > 0:
                time_data.append({
                    "userId": user_id,
                    "username": display_name,
                    "timeSpent": total_time
                })
        
        # Sort data
        streak_data.sort(key=lambda x: x["streak"], reverse=True)
        time_data.sort(key=lambda x: x["timeSpent"], reverse=True)
        
        # Limit to top 10
        top_streaks = streak_data[:10]
        top_learning_time = time_data[:10]
        
        # Update leaderboard in Firestore
        leaderboard_data = {
            "topStreaks": top_streaks,
            "topLearningTime": top_learning_time,
            "lastUpdated": datetime.now().isoformat()
        }
        
        db.collection("leaderboard").document("stats").set(leaderboard_data)
        
        return jsonify({
            "success": True,
            "message": "Leaderboard updated successfully",
            "leaderboard": leaderboard_data
        })
    except Exception as e:
        print(f"Error updating leaderboard: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Failed to update leaderboard: {str(e)}"}), 500

@app.route('/api/user/time-stats', methods=['POST'])
def get_user_time_stats():
    """Get user's time spent statistics"""
    try:
        data = request.json
        if not data:
            return jsonify({"error": "Invalid request: No JSON data provided"}), 400
            
        email = data.get("email")
        
        if not email:
            return jsonify({"error": "Email is required"}), 400
        
        # Find user in Firestore
        users_ref = db.collection("users")
        user_query = users_ref.where("email", "==", email).limit(1)
        user_docs = list(user_query.stream())
        
        if not user_docs:
            return jsonify({"error": "User not found"}), 404
        
        user_data = user_docs[0].to_dict()
        
        # Get stats or initialize if not present
        stats = user_data.get("stats", {})
        total_time = stats.get("totalTimeSpent", 0)
        activity_breakdown = stats.get("activityBreakdown", {})
        
        # Get current week time
        current_week_time = stats.get("currentWeekTime", 0)
        
        # Get last tracked week
        last_week_tracked = stats.get("lastWeekTracked", None)
        
        # Get the current week key
        today = date.today()
        year, week_num, _ = today.isocalendar()
        current_week_key = f"{year}-W{week_num:02d}"
        
        # Check if we need to update the lastWeekTracked
        is_new_week = last_week_tracked != current_week_key
        
        # Get recent sessions (last 7 days)
        sessions = stats.get("sessions", [])
        recent_sessions = []
        
        if sessions:
            now = datetime.now()
            one_week_ago = now - timedelta(days=7)
            
            for session in sessions:
                try:
                    session_time = datetime.fromisoformat(session.get("startTime", ""))
                    if session_time >= one_week_ago:
                        recent_sessions.append(session)
                except (ValueError, TypeError):
                    # If date format is invalid, skip this session
                    pass
        
        # Calculate daily averages
        daily_stats = {}
        for session in recent_sessions:
            try:
                session_date = datetime.fromisoformat(session.get("startTime", "")).date().isoformat()
                if session_date not in daily_stats:
                    daily_stats[session_date] = 0
                daily_stats[session_date] += session.get("durationMinutes", 0)
            except (ValueError, TypeError):
                # If date format is invalid, skip this session
                pass
        
        # Convert to list for easier frontend processing
        daily_time = [{"date": date, "minutes": minutes} for date, minutes in daily_stats.items()]
        
        return jsonify({
            "success": True,
            "totalTimeSpent": total_time,
            "currentWeekTime": current_week_time,
            "isNewWeek": is_new_week,
            "currentWeekKey": current_week_key,
            "lastWeekTracked": last_week_tracked,
            "activityBreakdown": activity_breakdown,
            "dailyTime": daily_time,
            "recentSessions": recent_sessions[-10:]  # Return last 10 sessions
        })
    except Exception as e:
        print(f"Error getting time stats: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Failed to get time stats: {str(e)}"}), 500

@app.route('/api/user/current-week-time', methods=['POST'])
def get_current_week_time():
    """Get user's current week learning time"""
    try:
        data = request.json
        if not data:
            return jsonify({"error": "Invalid request: No JSON data provided"}), 400
            
        email = data.get("email")
        
        if not email:
            return jsonify({"error": "Email is required"}), 400
        
        # Find user in Firestore
        users_ref = db.collection("users")
        user_query = users_ref.where("email", "==", email).limit(1)
        user_docs = list(user_query.stream())
        
        if not user_docs:
            return jsonify({"error": "User not found"}), 404
        
        user_data = user_docs[0].to_dict()
        
        # Get stats or initialize if not present
        stats = user_data.get("stats", {})
        
        # Get current week time
        current_week_time = stats.get("currentWeekTime", 0)
        
        # Get the current week from weekly stats if available
        weekly_stats = stats.get("weeklyStats", {})
        
        # Get today's date and current week key
        today = date.today()
        year, week_num, _ = today.isocalendar()
        current_week_key = f"{year}-W{week_num:02d}"
        
        # If current week exists in weekly stats, use that data
        if current_week_key in weekly_stats:
            current_week_time = weekly_stats[current_week_key].get("totalMinutes", 0)
        
        # Get last tracked week
        last_week_tracked = stats.get("lastWeekTracked", None)
        
        # Check if we need to update the lastWeekTracked
        is_new_week = last_week_tracked != current_week_key
        
        return jsonify({
            "success": True,
            "currentWeekTime": current_week_time,
            "isNewWeek": is_new_week,
            "currentWeekKey": current_week_key
        })
    except Exception as e:
        print(f"Error getting current week time: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Failed to get current week time: {str(e)}"}), 500

@app.route('/api/user/weekly-stats', methods=['POST'])
def get_weekly_stats():
    """Get user's weekly time statistics"""
    try:
        data = request.json
        if not data:
            return jsonify({"error": "Invalid request: No JSON data provided"}), 400
            
        email = data.get("email")
        
        if not email:
            return jsonify({"error": "Email is required"}), 400
        
        # Find user in Firestore
        users_ref = db.collection("users")
        user_query = users_ref.where("email", "==", email).limit(1)
        user_docs = list(user_query.stream())
        
        if not user_docs:
            return jsonify({"error": "User not found"}), 404
        
        user_data = user_docs[0].to_dict()
        
        # Get stats or initialize if not present
        stats = user_data.get("stats", {})
        total_time = stats.get("totalTimeSpent", 0)
        weekly_stats = stats.get("weeklyStats", {})
        
        # Get the current and past 5 weeks for display
        today = date.today()
        current_year, current_week, _ = today.isocalendar()
        
        recent_weeks = []
        for i in range(5, -1, -1):  # Last 6 weeks including current week
            # Calculate the date for i weeks ago
            past_date = today - timedelta(weeks=i)
            year, week_num, _ = past_date.isocalendar()
            week_key = f"{year}-W{week_num:02d}"
            
            week_data = weekly_stats.get(week_key, {
                "totalMinutes": 0,
                "activeDays": [],
                "activities": {}
            })
            
            # Calculate start and end dates for this week
            first_day = datetime.fromisocalendar(year, week_num, 1)
            last_day = datetime.fromisocalendar(year, week_num, 7)
            
            recent_weeks.append({
                "weekId": week_key,
                "startDate": first_day.strftime("%Y-%m-%d"),
                "endDate": last_day.strftime("%Y-%m-%d"),
                "hours": round(week_data.get("totalMinutes", 0) / 60, 1),
                "minutes": week_data.get("totalMinutes", 0),
                "activeDays": len(week_data.get("activeDays", [])),
                "activities": week_data.get("activities", {})
            })
        
        return jsonify({
            "success": True,
            "totalTimeSpent": total_time,
            "weeklyStats": recent_weeks,
            "hoursByWeek": [
                {"week": week["weekId"], "hours": week["hours"]} 
                for week in recent_weeks
            ]
        })
    except Exception as e:
        print(f"Error getting weekly stats: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": f"Failed to get weekly stats: {str(e)}"}), 500

@app.route('/api/user/learning-style', methods=['POST'])
def get_learning_style():
    """Get user's learning style data from Firestore"""
    try:
        data = request.json
        if not data:
            return jsonify({"error": "Invalid request: No JSON data provided"}), 400
            
        email = data.get("email")
        
        if not email:
            return jsonify({"error": "Email is required"}), 400
        
        # Find user in Firestore
        users_ref = db.collection("users")
        user_query = users_ref.where("email", "==", email).limit(1)
        user_docs = list(user_query.stream())
        
        if not user_docs:
            return jsonify({
                "success": False, 
                "error": "User not found",
                "primaryStyle": "Visual",
                "styles": [
                    {"name": "Visual", "percentage": 25},
                    {"name": "Auditory", "percentage": 25},
                    {"name": "Reading/Writing", "percentage": 25},
                    {"name": "Kinesthetic", "percentage": 25}
                ]
            }), 404
        
        user_data = user_docs[0].to_dict()
        
        # Get learning style data
        learning_style = user_data.get("learningStyle", "Visual")
        learning_style_details = user_data.get("learningStyleDetails", {})
        
        # Convert the details to a format suitable for the chart
        styles = []
        
        # Check the format of learning style details
        if isinstance(learning_style_details, dict) and learning_style_details:
            # If it's a dictionary with percentages
            total = sum(learning_style_details.values())
            
            for style, value in learning_style_details.items():
                percentage = int((value / total) * 100) if total > 0 else 25
                styles.append({
                    "name": style,
                    "percentage": percentage
                })
        else:
            # Use a default distribution based on primary style
            if learning_style == "Visual":
                styles = [
                    {"name": "Visual", "percentage": 40},
                    {"name": "Auditory", "percentage": 20},
                    {"name": "Reading/Writing", "percentage": 20},
                    {"name": "Kinesthetic", "percentage": 20}
                ]
            elif learning_style == "Auditory":
                styles = [
                    {"name": "Visual", "percentage": 20},
                    {"name": "Auditory", "percentage": 40},
                    {"name": "Reading/Writing", "percentage": 20},
                    {"name": "Kinesthetic", "percentage": 20}
                ]
            elif learning_style == "Reading/Writing":
                styles = [
                    {"name": "Visual", "percentage": 20},
                    {"name": "Auditory", "percentage": 20},
                    {"name": "Reading/Writing", "percentage": 40},
                    {"name": "Kinesthetic", "percentage": 20}
                ]
            elif learning_style == "Kinesthetic":
                styles = [
                    {"name": "Visual", "percentage": 20},
                    {"name": "Auditory", "percentage": 20},
                    {"name": "Reading/Writing", "percentage": 20},
                    {"name": "Kinesthetic", "percentage": 40}
                ]
            else:
                styles = [
                    {"name": "Visual", "percentage": 25},
                    {"name": "Auditory", "percentage": 25},
                    {"name": "Reading/Writing", "percentage": 25},
                    {"name": "Kinesthetic", "percentage": 25}
                ]
        
        return jsonify({
            "success": True,
            "primaryStyle": learning_style,
            "styles": styles,
            "lastUpdated": user_data.get("learningStyleUpdatedAt", None)
        })
        
    except Exception as e:
        print(f"Error getting learning style: {str(e)}")
        traceback.print_exc()
        return jsonify({
            "success": False, 
            "error": str(e),
            "primaryStyle": "Visual",
            "styles": [
                {"name": "Visual", "percentage": 25},
                {"name": "Auditory", "percentage": 25},
                {"name": "Reading/Writing", "percentage": 25},
                {"name": "Kinesthetic", "percentage": 25}
            ]
        }), 500

# User Preferences API routes
@app.route('/api/users/<user_id>/preferences/learning-style', methods=['GET'])
def get_user_learning_style(user_id):
    try:
        # Get user document from Firestore
        user_ref = db.collection('users').document(user_id)
        user_doc = user_ref.get()
        
        if not user_doc.exists:
            return jsonify({"error": "User not found"}), 404
            
        preferences_ref = user_ref.collection('preferences').document('learning_style')
        preferences_doc = preferences_ref.get()
        
        if not preferences_doc.exists:
            # Return default if not set
            return jsonify({
                "success": True,
                "defaultStyle": "visual"
            })
            
        preferences_data = preferences_doc.to_dict()
        return jsonify({
            "success": True,
            "defaultStyle": preferences_data.get("defaultStyle", "visual")
        })
    except Exception as e:
        print(f"Error getting user learning style: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/users/<user_id>/preferences/learning-style', methods=['POST'])
def set_user_learning_style(user_id):
    try:
        data = request.json
        if not data or 'defaultStyle' not in data:
            return jsonify({"error": "Default style is required"}), 400
            
        learning_style = data['defaultStyle']
        
        # Get user document from Firestore
        user_ref = db.collection('users').document(user_id)
        user_doc = user_ref.get()
        
        if not user_doc.exists:
            return jsonify({"error": "User not found"}), 404
            
        # Update or create preferences document
        preferences_ref = user_ref.collection('preferences').document('learning_style')
        preferences_ref.set({
            "defaultStyle": learning_style,
            "updatedAt": datetime.now().isoformat()
        }, merge=True)
        
        return jsonify({
            "success": True,
            "message": "Learning style preference updated"
        })
    except Exception as e:
        print(f"Error setting user learning style: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/users/<user_id>/preferences/subjects', methods=['GET'])
def get_user_subject_preferences(user_id):
    try:
        # Get user document from Firestore
        user_ref = db.collection('users').document(user_id)
        user_doc = user_ref.get()
        
        if not user_doc.exists:
            return jsonify({"error": "User not found"}), 404
            
        # Get all subject preferences
        subjects_ref = user_ref.collection('preferences').document('subjects')
        subjects_doc = subjects_ref.get()
        
        if not subjects_doc.exists:
            return jsonify({
                "success": True,
                "subjects": []
            })
            
        subjects_data = subjects_doc.to_dict()
        
        # Format the subjects into a list for the frontend
        subject_list = []
        for subject, style in subjects_data.get('subjectStyles', {}).items():
            subject_list.append({
                "subject": subject,
                "learningStyle": style
            })
            
        return jsonify({
            "success": True,
            "subjects": subject_list
        })
    except Exception as e:
        print(f"Error getting user subject preferences: {e}")
        return jsonify({"success": False, "error": str(e)}), 500
        
@app.route('/api/users/<user_id>/preferences/subjects', methods=['POST'])
def set_user_subject_preference(user_id):
    try:
        data = request.json
        if not data or 'subject' not in data or 'learningStyle' not in data:
            return jsonify({"error": "Subject and learningStyle are required"}), 400
            
        subject = data['subject']
        learning_style = data['learningStyle']
        
        # Get user document from Firestore
        user_ref = db.collection('users').document(user_id)
        user_doc = user_ref.get()
        
        if not user_doc.exists:
            return jsonify({"error": "User not found"}), 404
            
        # Update or create subject preferences
        subjects_ref = user_ref.collection('preferences').document('subjects')
        subjects_doc = subjects_ref.get()
        
        if subjects_doc.exists:
            subjects_data = subjects_doc.to_dict()
            subject_styles = subjects_data.get('subjectStyles', {})
            subject_styles[subject] = learning_style
            
            subjects_ref.update({
                'subjectStyles': subject_styles,
                'updatedAt': datetime.now().isoformat()
            })
        else:
            subjects_ref.set({
                'subjectStyles': {subject: learning_style},
                'updatedAt': datetime.now().isoformat()
            })
        
        return jsonify({
            "success": True,
            "message": "Subject preference updated"
        })
    except Exception as e:
        print(f"Error setting user subject preference: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/users/<user_id>/preferences/subjects/<subject>', methods=['DELETE'])
def delete_user_subject_preference(user_id, subject):
    try:
        # Get user document from Firestore
        user_ref = db.collection('users').document(user_id)
        user_doc = user_ref.get()
        
        if not user_doc.exists:
            return jsonify({"error": "User not found"}), 404
            
        # Get subject preferences
        subjects_ref = user_ref.collection('preferences').document('subjects')
        subjects_doc = subjects_ref.get()
        
        if not subjects_doc.exists:
            return jsonify({
                "success": True,
                "message": "No subjects found to delete"
            })
            
        subjects_data = subjects_doc.to_dict()
        subject_styles = subjects_data.get('subjectStyles', {})
        
        # Remove the subject if it exists
        if subject in subject_styles:
            subject_styles.pop(subject)
            
            # Update Firestore document
            subjects_ref.update({
                'subjectStyles': subject_styles,
                'updatedAt': datetime.now().isoformat()
            })
        
        return jsonify({
            "success": True,
            "message": "Subject preference deleted"
        })
    except Exception as e:
        print(f"Error deleting user subject preference: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/users/<user_id>/learning-effectiveness', methods=['GET'])
def get_user_learning_effectiveness(user_id):
    try:
        # Get user document from Firestore
        user_ref = db.collection('users').document(user_id)
        user_doc = user_ref.get()
        
        if not user_doc.exists:
            return jsonify({"error": "User not found"}), 404
            
        # Get learning effectiveness data
        effectiveness_ref = user_ref.collection('learning').document('effectiveness')
        effectiveness_doc = effectiveness_ref.get()
        
        if not effectiveness_doc.exists:
            return jsonify({
                "success": True,
                "styles": [],
                "recommendedStyle": None
            })
            
        effectiveness_data = effectiveness_doc.to_dict()
        
        # Calculate recommended style based on scores
        style_scores = effectiveness_data.get('styleScores', {})
        recommended_style = None
        best_score = 0
        
        for style, score_data in style_scores.items():
            avg_score = score_data.get('averageScore', 0)
            if avg_score > best_score:
                best_score = avg_score
                recommended_style = style
                
        # Format style data for frontend
        styles_list = []
        for style, score_data in style_scores.items():
            styles_list.append({
                "style": style,
                "averageScore": score_data.get('averageScore', 0),
                "quizCount": score_data.get('quizCount', 0)
            })
            
        return jsonify({
            "success": True,
            "styles": styles_list,
            "recommendedStyle": recommended_style
        })
    except Exception as e:
        print(f"Error getting learning effectiveness data: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/users/<user_id>/learning-effectiveness', methods=['POST'])
def track_learning_effectiveness(user_id):
    try:
        data = request.json
        if not data or 'learningStyle' not in data or 'score' not in data:
            return jsonify({"error": "Learning style and score are required"}), 400
            
        learning_style = data['learningStyle']
        score = float(data['score'])
        quiz_id = data.get('quizId', '')
        document_id = data.get('documentId', '')
        timestamp = data.get('timestamp', datetime.now().isoformat())
        
        # Get user document from Firestore
        user_ref = db.collection('users').document(user_id)
        user_doc = user_ref.get()
        
        if not user_doc.exists:
            return jsonify({"error": "User not found"}), 404
            
        # Update effectiveness document
        effectiveness_ref = user_ref.collection('learning').document('effectiveness')
        effectiveness_doc = effectiveness_ref.get()
        
        if effectiveness_doc.exists:
            effectiveness_data = effectiveness_doc.to_dict()
            style_scores = effectiveness_data.get('styleScores', {})
            
            # If this style doesn't exist yet, initialize it
            if learning_style not in style_scores:
                style_scores[learning_style] = {
                    'totalScore': 0,
                    'quizCount': 0,
                    'averageScore': 0
                }
                
            # Update the stats for this learning style
            style_data = style_scores[learning_style]
            new_total = style_data.get('totalScore', 0) + score
            new_count = style_data.get('quizCount', 0) + 1
            new_average = new_total / new_count
            
            style_scores[learning_style] = {
                'totalScore': new_total,
                'quizCount': new_count,
                'averageScore': new_average
            }
            
            # Update the document
            effectiveness_ref.update({
                'styleScores': style_scores,
                'updatedAt': datetime.now().isoformat()
            })
        else:
            # Create a new document
            effectiveness_ref.set({
                'styleScores': {
                    learning_style: {
                        'totalScore': score,
                        'quizCount': 1,
                        'averageScore': score
                    }
                },
                'createdAt': datetime.now().isoformat(),
                'updatedAt': datetime.now().isoformat()
            })
            
        # Also log this specific quiz result
        history_ref = user_ref.collection('learning').document('history')
        history_doc = history_ref.get()
        
        if history_doc.exists:
            history_data = history_doc.to_dict()
            quiz_history = history_data.get('quizHistory', [])
            quiz_history.append({
                'quizId': quiz_id,
                'documentId': document_id,
                'learningStyle': learning_style,
                'score': score,
                'timestamp': timestamp
            })
            
            # Limit history to last 50 quizzes
            if len(quiz_history) > 50:
                quiz_history = quiz_history[-50:]
                
            history_ref.update({
                'quizHistory': quiz_history,
                'updatedAt': datetime.now().isoformat()
            })
        else:
            history_ref.set({
                'quizHistory': [{
                    'quizId': quiz_id,
                    'documentId': document_id,
                    'learningStyle': learning_style,
                    'score': score,
                    'timestamp': timestamp
                }],
                'createdAt': datetime.now().isoformat(),
                'updatedAt': datetime.now().isoformat()
            })
        
        return jsonify({
            "success": True,
            "message": "Learning effectiveness data updated"
        })
    except Exception as e:
        print(f"Error tracking learning effectiveness: {e}")
        return jsonify({"success": False, "error": str(e)}), 500

if __name__ == "__main__":
    # Set debug mode to False in production to avoid auto-reloading
    # You can enable it with an environment variable for development
    debug_mode = os.getenv('FLASK_DEBUG', 'False').lower() == 'true'
    app.run(debug=debug_mode, host='0.0.0.0', port=int(os.getenv('PORT', 5000)))
